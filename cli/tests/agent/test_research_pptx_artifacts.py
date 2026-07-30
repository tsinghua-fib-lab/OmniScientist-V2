"""research-pptx artifact contract and presentation integration."""

from __future__ import annotations

import importlib.util
import json
import os
import sys
import time
from pathlib import Path
from types import SimpleNamespace

import pytest
import yaml

from omni.channels.outbound import delivery_envelope_from_presentation
from omni.config import load_settings
from omni.runtime.presentation import task_presentation_from_result
from omni.storage.artifacts import ArtifactStore
from omni.storage.db import get_database

SKILL_DIR = Path(__file__).resolve().parents[3] / "skills" / "research-pptx"
PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"


def _load_engine():
    path = SKILL_DIR / "engine.py"
    spec = importlib.util.spec_from_file_location("test_research_pptx_engine", path)
    assert spec is not None and spec.loader is not None
    module = importlib.util.module_from_spec(spec)
    sys.modules[spec.name] = module
    spec.loader.exec_module(module)
    return module


def _load_portable_runner():
    path = SKILL_DIR / "scripts" / "run.py"
    spec = importlib.util.spec_from_file_location("test_research_pptx_runner", path)
    assert spec is not None and spec.loader is not None
    module = importlib.util.module_from_spec(spec)
    sys.modules[spec.name] = module
    spec.loader.exec_module(module)
    return module


def _load_overflow_script():
    path = SKILL_DIR / "scripts" / "check_overflow.py"
    spec = importlib.util.spec_from_file_location("test_research_pptx_overflow", path)
    assert spec is not None and spec.loader is not None
    module = importlib.util.module_from_spec(spec)
    sys.modules[spec.name] = module
    spec.loader.exec_module(module)
    return module


@pytest.mark.asyncio
async def test_research_pptx_store_returns_complete_artifact_descriptor(tmp_path: Path) -> None:
    settings = load_settings()
    settings.paths.ensure_dirs()
    db = get_database(settings.paths.project_db)
    await db.init()
    store = ArtifactStore(settings.paths, db)
    engine = _load_engine().ResearchPptxEngine()
    engine.ctx = SimpleNamespace(
        artifacts=store,
        session_id="",
        subtask_id="",
        workflow_run_id="",
    )
    source = tmp_path / "deck.pptx"
    source.write_bytes(b"pptx-bytes")

    artifact = await engine._store_pptx(str(source), "RAG Group Meeting")  # noqa: SLF001

    assert artifact["title"] == "RAG Group Meeting"
    assert artifact["format"] == "pptx"
    assert artifact["uri"].startswith("artifact://")
    assert Path(artifact["path"]).is_file()
    assert artifact["mime"] == PPTX_MIME
    assert artifact["size_bytes"] == len(b"pptx-bytes")


@pytest.mark.asyncio
async def test_research_pptx_store_passes_execution_ownership(tmp_path: Path) -> None:
    source = tmp_path / "deck.pptx"
    source.write_bytes(b"pptx-bytes")
    captured: dict = {}

    class CapturingArtifacts:
        async def put_bytes(self, data: bytes, **kwargs):  # noqa: ANN003, ANN202
            captured.update(kwargs)
            return SimpleNamespace(
                uri="artifact://deck123",
                path=source,
                mime=PPTX_MIME,
                size_bytes=len(data),
            )

    engine = _load_engine().ResearchPptxEngine()
    engine.ctx = SimpleNamespace(
        artifacts=CapturingArtifacts(),
        session_id="session-pptx",
        subtask_id="subtask-pptx",
        workflow_run_id="workflow-pptx",
    )

    await engine._store_pptx(str(source), "RAG Group Meeting")  # noqa: SLF001

    assert captured["session_id"] == "session-pptx"
    assert captured["subtask_id"] == "subtask-pptx"
    assert captured["workflow_run_id"] == "workflow-pptx"


@pytest.mark.asyncio
async def test_portable_artifact_store_matches_engine_contract(tmp_path: Path) -> None:
    runner = _load_portable_runner()
    source = tmp_path / "deck.pptx"
    source.write_bytes(b"pptx-bytes")
    engine = _load_engine().ResearchPptxEngine()
    engine.ctx = SimpleNamespace(
        artifacts=runner._MockArtifacts(tmp_path / "portable-out"),  # noqa: SLF001
        session_id="portable-session",
        subtask_id="portable-subtask",
        workflow_run_id="portable-workflow",
    )

    artifact = await engine._store_pptx(str(source), "Portable Deck")  # noqa: SLF001

    assert artifact["uri"].startswith("file://")
    assert Path(artifact["path"]).read_bytes() == b"pptx-bytes"
    assert artifact["mime"] == PPTX_MIME
    assert artifact["size_bytes"] == len(b"pptx-bytes")


def test_research_pptx_manifest_declares_structured_artifacts() -> None:
    frontmatter = yaml.safe_load((SKILL_DIR / "SKILL.md").read_text(encoding="utf-8").split("---", 2)[1])
    helix = frontmatter["metadata"]["helixforge"]
    artifacts = helix["output_schema"]["properties"]["artifacts"]

    assert artifacts["type"] == "array"
    assert artifacts["items"]["required"] == ["title", "format", "uri", "path", "mime", "size_bytes"]
    assert helix["quality_contract"]["assessment_required"] is True
    assert helix["input_schema"]["properties"]["plan_edits"]["type"] == "array"


def test_research_pptx_result_model_preserves_artifacts() -> None:
    result = _load_engine().PresentationResult(
        title="RAG Group Meeting",
        pptx_uri="artifact://deck123",
        artifacts=[
            {
                "title": "RAG Group Meeting",
                "format": "pptx",
                "uri": "artifact://deck123",
                "path": "/workspace/artifacts/presentation/deck.pptx",
                "mime": PPTX_MIME,
                "size_bytes": 42,
            }
        ],
    ).model_dump()

    assert result["artifacts"][0]["uri"] == "artifact://deck123"


def test_research_pptx_request_remaps_common_source_and_count_aliases(
    tmp_path: Path,
) -> None:
    module = _load_engine()
    paper = tmp_path / "attention.pdf"
    paper.write_bytes(b"%PDF-test")

    request = module.PresentationRequest(
        paper_path=str(paper),
        slide_count=15,
        outline_path="/notes/talk.md",
    )

    assert request.pdf_uri == str(paper)
    assert request.target_slides == 15
    assert request.markdown_uri == "/notes/talk.md"

    arguments = {"paper_path": str(paper), "slide_count": 15}
    assert module.ResearchPptxEngine.validate_params(arguments=arguments) is None
    assert arguments["pdf_uri"] == str(paper)
    assert arguments["target_slides"] == 15


@pytest.mark.parametrize(
    ("field", "value"),
    [
        ("paper_uris", ["/papers/one.pdf", "/papers/two.pdf"]),
        ("file_uris", ["/notes/context.md"]),
        ("corpus_query", "retrieval augmented generation"),
        ("source_ids", ["source-1"]),
    ],
)
def test_research_pptx_validation_accepts_every_declared_source_kind(
    field: str,
    value: object,
) -> None:
    module = _load_engine()
    arguments = {field: value}

    assert module.ResearchPptxEngine.validate_params(arguments=arguments) is None


def test_research_pptx_does_not_treat_ambiguous_source_as_a_pdf() -> None:
    module = _load_engine()
    arguments = {"topic": "RAG", "source": "inline evidence from the discussion"}

    assert module.ResearchPptxEngine.validate_params(arguments=arguments) is None
    assert "pdf_uri" not in arguments
    assert module.PresentationRequest(**arguments).pdf_uri is None


def test_research_pptx_plan_edits_follow_declared_operation_order() -> None:
    module = _load_engine()
    plan = {
        "title": "Deck",
        "slides": [
            {"slide_type": "title", "title": "Title", "bullets": []},
            {"slide_type": "content", "title": "A", "bullets": ["a"]},
            {"slide_type": "content", "title": "B", "bullets": ["b"]},
            {"slide_type": "conclusion", "title": "End", "bullets": []},
        ],
    }

    edited = module.ResearchPptxEngine._apply_plan_edits(  # noqa: SLF001
        plan,
        [
            {"action": "set_title", "slide_index": 2, "title": "B revised"},
            {"action": "remove_slide", "slide_index": 1},
        ],
    )

    assert [slide["title"] for slide in edited["slides"]] == [
        "Title",
        "B revised",
        "End",
    ]


def test_research_pptx_plan_edits_are_atomic_and_fail_closed() -> None:
    module = _load_engine()
    plan = {
        "title": "Deck",
        "slides": [
            {"slide_type": "title", "title": "Title", "bullets": []},
            {"slide_type": "content", "title": "Evidence", "bullets": ["old"]},
            {"slide_type": "conclusion", "title": "End", "bullets": []},
        ],
    }

    edited = module.ResearchPptxEngine._apply_plan_edits(  # noqa: SLF001
        plan,
        [{"action": "set_bullet", "slide_index": 1, "bullet_index": 0, "text": "new"}],
    )

    assert edited["slides"][1]["bullets"] == ["new"]
    assert plan["slides"][1]["bullets"] == ["old"]
    with pytest.raises(ValueError, match="unsupported plan edit action"):
        module.ResearchPptxEngine._apply_plan_edits(  # noqa: SLF001
            plan,
            [{"action": "invent_slide", "slide_index": 1}],
        )
    with pytest.raises(ValueError, match="slide_index"):
        module.ResearchPptxEngine._apply_plan_edits(  # noqa: SLF001
            plan,
            [{"action": "set_title", "slide_index": 99, "title": "missing"}],
        )


def test_research_pptx_auto_resume_is_scoped_and_unambiguous(tmp_path: Path) -> None:
    module = _load_engine()
    review_root = tmp_path / "pptx_review"

    def write_state(token: str, *, session_id: str, age_seconds: int) -> None:
        state_dir = review_root / token
        state_dir.mkdir(parents=True)
        state_file = state_dir / "state.json"
        state_file.write_text(
            json.dumps(
                {
                    "plan": {"title": token, "slides": []},
                    "content": {},
                    "owner": {"session_id": session_id, "task_id": f"task-{token}"},
                }
            ),
            encoding="utf-8",
        )
        stamp = time.time() - age_seconds
        os.utime(state_file, (stamp, stamp))

    write_state("owned", session_id="session-a", age_seconds=20)
    write_state("other", session_id="session-b", age_seconds=1)
    engine = module.ResearchPptxEngine()
    engine.ctx = SimpleNamespace(
        paths=SimpleNamespace(artifacts_dir=tmp_path),
        session_id="session-a",
        task_id="new-task",
    )

    assert engine._find_recent_review_token() == "owned"  # noqa: SLF001
    approval = {"topic": "approve"}
    assert engine._detect_misrouted_intent(approval) is None  # noqa: SLF001
    assert approval["resume_token"] == "owned"

    write_state("also-owned", session_id="session-a", age_seconds=10)
    assert engine._find_recent_review_token() is None  # noqa: SLF001


@pytest.mark.parametrize("candidate_count", [0, 2])
def test_research_pptx_auto_resume_without_one_candidate_fails_closed(
    tmp_path: Path,
    candidate_count: int,
) -> None:
    module = _load_engine()
    review_root = tmp_path / "pptx_review"
    for index in range(candidate_count):
        state_dir = review_root / f"token-{index}"
        state_dir.mkdir(parents=True)
        (state_dir / "state.json").write_text(
            json.dumps(
                {
                    "plan": {"title": "Deck", "slides": []},
                    "content": {"source_type": "prompt"},
                    "owner": {"session_id": "session-a"},
                    "consumed_at": None,
                }
            ),
            encoding="utf-8",
        )
    engine = module.ResearchPptxEngine()
    engine.ctx = SimpleNamespace(
        paths=SimpleNamespace(artifacts_dir=tmp_path),
        session_id="session-a",
    )
    arguments = {"topic": "approve the outline"}

    result = engine._detect_misrouted_intent(arguments)  # noqa: SLF001

    assert result is not None
    assert result["status"] == "error"
    assert result["outcome"] == {"code": "review_pending"}
    assert result["recoverable"] is True
    assert "resume_token" in result["error"]
    assert arguments == {"topic": "approve the outline"}


@pytest.mark.parametrize(
    ("topic", "expected_topic"),
    [
        ("make a 15-slide deck; show the outline first", "make a 15-slide deck"),
        ("制作15页PPT，先看大纲", "制作15页PPT"),
    ],
)
def test_research_pptx_extracts_page_count_before_outline_review_routing(
    topic: str,
    expected_topic: str,
) -> None:
    module = _load_engine()
    engine = module.ResearchPptxEngine()
    engine.ctx = SimpleNamespace(paths=None, session_id="session-a")
    arguments = {"topic": topic}

    assert engine._detect_misrouted_intent(arguments) is None  # noqa: SLF001
    assert arguments == {
        "topic": expected_topic,
        "review_mode": "plan",
        "target_slides": 15,
    }


@pytest.mark.asyncio
async def test_research_pptx_review_state_persists_original_request(tmp_path: Path) -> None:
    module = _load_engine()
    engine = module.ResearchPptxEngine()
    engine.ctx = SimpleNamespace(
        paths=SimpleNamespace(artifacts_dir=tmp_path),
        session_id="session-a",
        task_id="task-a",
        subtask_id="",
        workflow_run_id="",
    )
    plan = module.PresentationPlan(
        title="Deck",
        slides=[
            module._models.SlideData(slide_type="title", title="Deck"),  # noqa: SLF001
            module._models.SlideData(  # noqa: SLF001
                slide_type="content", title="Evidence", bullets=["A", "B"]
            ),
            module._models.SlideData(slide_type="conclusion", title="End"),  # noqa: SLF001
        ],
    )
    request = module.PresentationRequest(
        topic="RAG",
        language="zh",
        target_slides=3,
        review_mode="plan",
    )

    await engine._persist_review_state(  # noqa: SLF001
        "review-token",
        plan,
        module.ParsedContent(source_type="prompt", markdown_text="RAG"),
        str(tmp_path / "work"),
        request=request,
    )

    state = json.loads(
        (tmp_path / "pptx_review" / "review-token" / "state.json").read_text(
            encoding="utf-8"
        )
    )
    assert state["request"]["language"] == "zh"
    assert state["request"]["target_slides"] == 3
    assert state["request"]["review_mode"] == "plan"


@pytest.mark.asyncio
async def test_research_pptx_resume_rejects_invalid_edits_without_consuming_token(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    module = _load_engine()
    token = "invalid-edit-token"
    monkeypatch.setitem(
        module._REVIEW_CACHE,  # noqa: SLF001
        token,
        {
            "plan": {
                "title": "Deck",
                "slides": [
                    {"slide_type": "title", "title": "Deck"},
                    {"slide_type": "content", "title": "Evidence", "bullets": ["A"]},
                    {"slide_type": "conclusion", "title": "End"},
                ],
            },
            "content": {"source_type": "prompt", "markdown_text": "RAG"},
            "request": {"topic": "RAG", "language": "zh", "target_slides": 3},
        },
    )
    engine = module.ResearchPptxEngine()
    engine.ctx = SimpleNamespace(paths=None, session_id="session-a")

    async def fail_if_rendered(*_args: object, **_kwargs: object) -> object:
        pytest.fail("invalid plan_edits must not reach rendering")

    monkeypatch.setattr(engine, "_render_and_finish", fail_if_rendered)

    result = await engine._run(  # noqa: SLF001
        module.PresentationRequest(
            resume_token=token,
            plan_edits=[{"action": "remove_slide", "slide_index": 99}],
        ),
        None,
    )

    assert result["status"] == "error"
    assert result["outcome"] == {"code": "invalid_plan_edits"}
    assert token in module._REVIEW_CACHE  # noqa: SLF001


@pytest.mark.asyncio
async def test_research_pptx_resume_restores_and_enforces_exact_slide_count(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    module = _load_engine()
    token = "count-token"
    monkeypatch.setitem(
        module._REVIEW_CACHE,  # noqa: SLF001
        token,
        {
            "plan": {
                "title": "Deck",
                "slides": [
                    {"slide_type": "title", "title": "Deck"},
                    {"slide_type": "content", "title": "Evidence", "bullets": ["A"]},
                    {"slide_type": "conclusion", "title": "End"},
                ],
            },
            "content": {"source_type": "prompt", "markdown_text": "RAG"},
            "request": {"topic": "RAG", "language": "zh", "target_slides": 3},
        },
    )
    engine = module.ResearchPptxEngine()
    engine.ctx = SimpleNamespace(paths=None, session_id="session-a")

    async def fail_if_rendered(*_args: object, **_kwargs: object) -> object:
        pytest.fail("a count-mismatched reviewed plan must not render")

    monkeypatch.setattr(engine, "_render_and_finish", fail_if_rendered)

    result = await engine._run(  # noqa: SLF001
        module.PresentationRequest(
            resume_token=token,
            plan_edits=[{"action": "remove_slide", "slide_index": 1}],
        ),
        None,
    )

    assert result["status"] == "error"
    assert result["outcome"] == {"code": "slide_count_mismatch"}
    assert result["target_slides"] == 3
    assert result["actual_slides"] == 2
    assert token in module._REVIEW_CACHE  # noqa: SLF001


@pytest.mark.asyncio
async def test_research_pptx_resume_restores_original_render_request(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    module = _load_engine()
    token = "restored-request-token"
    monkeypatch.setitem(
        module._REVIEW_CACHE,  # noqa: SLF001
        token,
        {
            "plan": {
                "title": "Deck",
                "slides": [
                    {"slide_type": "title", "title": "Deck"},
                    {"slide_type": "content", "title": "Evidence", "bullets": ["A"]},
                    {"slide_type": "conclusion", "title": "End"},
                ],
            },
            "content": {"source_type": "prompt", "markdown_text": "RAG"},
            "request": {
                "topic": "RAG",
                "language": "zh",
                "talk_type": "defense",
                "target_slides": 3,
            },
        },
    )
    engine = module.ResearchPptxEngine()
    engine.ctx = SimpleNamespace(paths=None, session_id="session-a")
    captured: dict[str, object] = {}

    async def capture_render(
        _plan: object,
        _work_dir: str,
        request: object,
        *_args: object,
    ) -> object:
        captured["request"] = request
        return module.PresentationResult(
            status="ok",
            title="Deck",
            slide_count=3,
        )

    monkeypatch.setattr(engine, "_render_and_finish", capture_render)

    result = await engine._run(  # noqa: SLF001
        module.PresentationRequest(resume_token=token),
        None,
    )

    restored = captured["request"]
    assert restored.language == "zh"
    assert restored.talk_type == "defense"
    assert restored.target_slides == 3
    assert result.status == "ok"
    assert token not in module._REVIEW_CACHE  # noqa: SLF001


def test_research_pptx_sparse_fix_preserves_explicit_page_count() -> None:
    module = _load_engine()
    plan = module.PresentationPlan(
        title="Deck",
        slides=[
            module._models.SlideData(slide_type="title", title="Deck"),  # noqa: SLF001
            module._models.SlideData(slide_type="content", title="Sparse"),  # noqa: SLF001
            module._models.SlideData(slide_type="conclusion", title="End"),  # noqa: SLF001
        ],
    )
    warnings = module.ResearchPptxEngine._check_sparse_slides(plan)  # noqa: SLF001

    module.ResearchPptxEngine._fix_sparse_slides(  # noqa: SLF001
        plan,
        warnings,
        target_slides=3,
    )

    assert len(plan.slides) == 3


@pytest.mark.asyncio
async def test_research_pptx_render_fails_closed_on_exact_count_mismatch() -> None:
    module = _load_engine()
    engine = module.ResearchPptxEngine()
    engine.ctx = SimpleNamespace()
    plan = module.PresentationPlan(
        title="Deck",
        slides=[
            module._models.SlideData(slide_type="title", title="Deck"),  # noqa: SLF001
            module._models.SlideData(slide_type="conclusion", title="End"),  # noqa: SLF001
        ],
    )

    result = await engine._render_and_finish(  # noqa: SLF001
        plan,
        str(Path.cwd()),
        module.PresentationRequest(topic="RAG", target_slides=3),
        SimpleNamespace(),
        None,
        time.time(),
        lambda: 0,
    )

    assert result["status"] == "error"
    assert result["outcome"] == {"code": "slide_count_mismatch"}
    assert result["phase"] == "pre_render"


@pytest.mark.asyncio
async def test_research_pptx_expands_an_explicit_slide_count_exactly() -> None:
    module = _load_engine()
    planner = module._load_sibling("content_planner")  # noqa: SLF001

    class FakeLLM:
        async def chat(self, *_args: object, **_kwargs: object) -> str:
            return json.dumps(
                {
                    "title": "RAG",
                    "slides": [
                        {"slide_type": "title", "title": "RAG"},
                        {
                            "slide_type": "content",
                            "title": "Retrieval grounds generation",
                            "bullets": [
                                "Query encoding",
                                "Dense retrieval",
                                "Reranking",
                                "Generation",
                            ],
                        },
                        {"slide_type": "conclusion", "title": "Takeaways"},
                    ],
                }
            )

    request = module.PresentationRequest(topic="RAG", target_slides=5)
    content = module.ParsedContent(source_type="prompt", markdown_text="RAG systems")

    plan = await planner.plan_presentation(FakeLLM(), content, request)

    assert len(plan.slides) == 5
    assert plan.slides[0].slide_type == "title"
    assert plan.slides[-1].slide_type == "conclusion"


def test_research_pptx_fits_over_and_exact_slide_counts() -> None:
    module = _load_engine()
    planner = module._load_sibling("content_planner")  # noqa: SLF001

    def make_plan(count: int):  # noqa: ANN202
        slides = [module._models.SlideData(slide_type="title", title="Title")]  # noqa: SLF001
        slides.extend(
            module._models.SlideData(  # noqa: SLF001
                slide_type="content",
                title=f"Evidence {index}",
                bullets=["A", "B"],
            )
            for index in range(count - 2)
        )
        slides.append(module._models.SlideData(slide_type="conclusion", title="End"))  # noqa: SLF001
        return module.PresentationPlan(title="Deck", slides=slides)

    over = planner._trim_slides_to_target(make_plan(7), 5)  # noqa: SLF001
    exact = planner._trim_slides_to_target(make_plan(5), 5)  # noqa: SLF001

    assert len(over.slides) == 5
    assert len(exact.slides) == 5


def test_research_pptx_trim_preserves_one_of_each_essential_slide_type() -> None:
    module = _load_engine()
    planner = module._load_sibling("content_planner")  # noqa: SLF001
    slide = module._models.SlideData  # noqa: SLF001
    plan = module.PresentationPlan(
        title="Evidence deck",
        slides=[
            slide(slide_type="title", title="Title"),
            slide(slide_type="table", title="Table A", table_headers=["x"], table_rows=[["1"]]),
            slide(slide_type="table", title="Table B", table_headers=["x"], table_rows=[["2"]]),
            slide(slide_type="metrics", title="Metrics A", metrics=[{"value": "1", "label": "A"}]),
            slide(slide_type="metrics", title="Metrics B", metrics=[{"value": "2", "label": "B"}]),
            slide(slide_type="full_figure", title="Figure A", figure_path="figure_0"),
            slide(slide_type="full_figure", title="Figure B", figure_path="figure_1"),
            slide(slide_type="conclusion", title="End"),
        ],
    )

    trimmed = planner._trim_slides_to_target(plan, 5)  # noqa: SLF001

    assert len(trimmed.slides) == 5
    remaining_types = [item.slide_type for item in trimmed.slides]
    assert remaining_types.count("table") == 1
    assert remaining_types.count("metrics") == 1
    assert remaining_types.count("full_figure") == 1


def test_research_pptx_table_insertion_keeps_highlight_inside_truncated_rows() -> None:
    module = _load_engine()
    planner = module._load_sibling("content_planner")  # noqa: SLF001
    slide = module._models.SlideData  # noqa: SLF001
    plan = module.PresentationPlan(
        title="Evidence deck",
        slides=[
            slide(slide_type="title", title="Title"),
            slide(slide_type="content", title="Intro", bullets=["A", "B"]),
            slide(slide_type="content", title="Results", bullets=["C", "D"]),
            slide(slide_type="conclusion", title="End"),
        ],
    )
    content = module.ParsedContent(
        source_type="prompt",
        markdown_text="Evidence",
        tables=[
            {
                "caption": "Results",
                "headers": ["Metric", "Value"],
                "rows": [[f"M{index}", str(index)] for index in range(8)],
            }
        ],
    )

    updated = planner._enforce_visual_variety(  # noqa: SLF001
        plan,
        content,
        target_slides=4,
    )

    table_slide = next(item for item in updated.slides if item.slide_type == "table")
    assert len(table_slide.table_rows) == 6
    assert table_slide.highlight_row == 5


def test_research_pptx_exact_target_references_never_drop_entries() -> None:
    module = _load_engine()
    planner = module._load_sibling("content_planner")  # noqa: SLF001
    slide = module._models.SlideData  # noqa: SLF001
    references = [
        {"key": f"[{index}]", "text": f"Reference {index}"}
        for index in range(1, 21)
    ]
    plan = module.PresentationPlan(
        title="Evidence deck",
        references=references,
        slides=[
            slide(slide_type="title", title="Title"),
            slide(slide_type="content", title="Intro", bullets=["A", "B"]),
            slide(slide_type="content", title="Evidence A", bullets=["C", "D"]),
            slide(slide_type="content", title="Evidence B", bullets=["E", "F"]),
            slide(slide_type="conclusion", title="End"),
        ],
    )

    updated = planner._append_references_slide(  # noqa: SLF001
        plan,
        target_slides=5,
    )

    reference_slides = [
        item for item in updated.slides if item.slide_type == "references"
    ]
    assert len(updated.slides) == 5
    assert len(reference_slides) == 2
    assert [bullet for item in reference_slides for bullet in item.bullets] == [
        f"{reference['key']} {reference['text']}" for reference in references
    ]


def test_research_pptx_exact_target_references_fail_when_no_slide_can_be_replaced() -> None:
    module = _load_engine()
    planner = module._load_sibling("content_planner")  # noqa: SLF001
    slide = module._models.SlideData  # noqa: SLF001
    plan = module.PresentationPlan(
        title="Minimal deck",
        references=[{"key": "[1]", "text": "Reference 1"}],
        slides=[
            slide(slide_type="title", title="Title"),
            slide(slide_type="content", title="Only evidence", bullets=["A", "B"]),
            slide(slide_type="conclusion", title="End"),
        ],
    )

    with pytest.raises(ValueError, match="reference slide"):
        planner._append_references_slide(plan, target_slides=3)  # noqa: SLF001


def test_research_pptx_template_bullet_xml_follows_drawingml_order() -> None:
    from pptx import Presentation
    from pptx.util import Inches

    module = _load_engine()
    template_backend = module._load_sibling("template_backend")  # noqa: SLF001
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    paragraph = slide.shapes.add_textbox(
        Inches(1), Inches(1), Inches(4), Inches(1)
    ).text_frame.paragraphs[0]
    paragraph.text = "Native bullet"
    paragraph.font.name = "Arial"

    template_backend._set_paragraph_bullet(paragraph)  # noqa: SLF001

    tags = [child.tag.rsplit("}", 1)[-1] for child in paragraph._p.pPr]  # noqa: SLF001
    assert tags.index("buFont") < tags.index("buChar") < tags.index("defRPr")


def test_research_pptx_table_overflow_uses_actual_narrow_column_widths() -> None:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    module = _load_engine()
    slide_renderer = module._load_sibling("slide_renderer")  # noqa: SLF001
    overflow_script = _load_overflow_script()
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    table_shape = slide.shapes.add_table(
        2,
        8,
        Inches(0.5),
        Inches(0.5),
        Inches(8.3),
        Inches(0.8),
    )
    for row in table_shape.table.rows:
        row.height = Inches(0.4)
        for cell in row.cells:
            cell.text = "检索增强生成系统架构与长文本内容测试"
            for run in cell.text_frame.paragraphs[0].runs:
                run.font.size = Pt(9)

    structured = slide_renderer._check_table_overflow_structured(  # noqa: SLF001
        table_shape,
        16,
        0.05,
    )
    standalone = overflow_script.check_table_overflow(table_shape, 1)

    assert structured > 0.3
    assert any("table" in warning for warning in standalone)


def test_research_pptx_template_sanitizes_every_display_text_field() -> None:
    module = _load_engine()
    template_backend = module._load_sibling("template_backend")  # noqa: SLF001
    slide = module._models.SlideData(  # noqa: SLF001
        slide_type="two_column",
        title="Title &bull;",
        subtitle="Subtitle &#8226;",
        bullets=["Bullet &bull;"],
        figure_path="/tmp/figure&amp;.png",
        figure_caption="Caption &bull;",
        notes="Notes &bull;",
        metrics=[{"value": "95&plusmn;2", "label": "Metric &bull;"}],
        table_headers=["Header &bull;"],
        table_rows=[["Cell &bull;"]],
        citations=[{"key": "[1]", "text": "Citation &bull;"}],
        extra={
            "columns": [
                {
                    "sub_title": "Column &bull;",
                    "bullets": ["Nested &bull;"],
                    "figure_caption": "Nested caption &bull;",
                    "figure_path": "/tmp/nested&amp;.png",
                }
            ],
            "rows": [
                {
                    "label": "Row &bull;",
                    "header": "Row header &bull;",
                    "description": "Row description &bull;",
                }
            ],
            "steps": [
                {
                    "step_number": "1&bull;",
                    "step_title": "Step &bull;",
                    "step_desc": "Step description &bull;",
                }
            ],
            "box_text": "Box &bull;",
            "emphasis_note": "Emphasis &bull;",
            "authors": "Author &bull;",
            "affiliation": "Lab &bull;",
        },
    )
    plan = module.PresentationPlan(  # noqa: SLF001
        title="Plan &bull;",
        authors="Authors &bull;",
        affiliation="Affiliation &bull;",
        venue="Venue &bull;",
        references=[{"key": "[1]", "text": "Reference &bull;"}],
        slides=[slide],
    )

    template_backend._sanitize_plan_entities(plan)  # noqa: SLF001

    rendered_text = json.dumps(plan.model_dump(), ensure_ascii=False)
    assert "&bull;" not in rendered_text
    assert "&#8226;" not in rendered_text
    assert "95±2" in rendered_text
    assert slide.figure_path == "/tmp/figure&amp;.png"
    assert slide.extra["columns"][0]["figure_path"] == "/tmp/nested&amp;.png"


@pytest.mark.asyncio
async def test_research_pptx_review_checkpoint_is_visible_and_nonblocking(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    module = _load_engine()
    planner = module._load_sibling("content_planner")  # noqa: SLF001
    plan = module.PresentationPlan(
        title="RAG Review",
        slides=[
            module._models.SlideData(slide_type="title", title="RAG Review"),  # noqa: SLF001
            module._models.SlideData(  # noqa: SLF001
                slide_type="content",
                title="Retrieval grounds generation",
                bullets=["Retrieve", "Rerank", "Generate"],
            ),
            module._models.SlideData(slide_type="conclusion", title="Takeaways"),  # noqa: SLF001
        ],
    )

    async def fake_process(*_args: object) -> object:
        return module.ParsedContent(source_type="prompt", markdown_text="RAG")

    async def fake_plan(*_args: object, **_kwargs: object) -> object:
        return plan

    monkeypatch.setattr(planner, "plan_presentation", fake_plan)
    engine = module.ResearchPptxEngine()
    monkeypatch.setattr(engine, "_process_inputs", fake_process)
    engine.ctx = SimpleNamespace(
        artifacts=None,
        llm=object(),
        paths=None,
        session_id="session-review",
        task_id="task-review",
    )
    progress: list[tuple[str, float]] = []

    async def record_progress(stage: str, fraction: float) -> None:
        progress.append((stage, fraction))

    result = await engine._run(  # noqa: SLF001
        module.PresentationRequest(topic="RAG", review_mode="plan"),
        record_progress,
    )

    assert result["status"] == "partial"
    assert result["outcome"] == {"code": "awaiting_review"}
    assert result["blocking"] is False
    assert result["_omni_control"]["terminal"] is True
    assert "# Presentation outline for review: RAG Review" in result["summary"]
    assert {stage for stage, _fraction in progress} <= {
        "parsing",
        "deciding",
        "planning",
        "rendering",
        "qa",
        "critique",
        "upload",
    }


@pytest.mark.asyncio
async def test_research_pptx_success_dict_keeps_deliverable_assessment(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    module = _load_engine()

    async def fake_pipeline(*_args: object, **_kwargs: object) -> dict[str, object]:
        return {
            "status": "ok",
            "pptx_uri": "artifact://deck123",
            "slide_count": 5,
            "artifacts": [{"uri": "artifact://deck123"}],
        }

    monkeypatch.setattr(module.ResearchPptxEngine, "_run", fake_pipeline)
    engine = module.ResearchPptxEngine()
    engine.ctx = SimpleNamespace()

    result = await engine.execute(topic="RAG")

    assert result["deliverable_assessment"]["status"] == "passed"


def test_research_pptx_python_dependency_error_returns_setup_action() -> None:
    module = _load_engine()
    python_result = module._python_dependency_error(  # noqa: SLF001
        ModuleNotFoundError("No module named 'fitz'", name="fitz")
    )
    assert python_result["error_info"]["code"] == "runtime_dependency_missing"
    assert python_result["action_required"]["kind"] == "install"
    assert python_result["setup_command"] == "omni update --force"
    assert python_result["action_required"]["missing"] == ["fitz"]


def test_pptx_presentation_names_the_deck_and_says_where_it_is(tmp_path: Path) -> None:
    pptx_path = tmp_path / "rag-group-meeting.pptx"
    pptx_path.write_bytes(b"pptx")
    result = {
        "status": "ok",
        "summary": "Generated the deck.",
        "pptx_uri": "artifact://deck123",
        "artifacts": [
            {
                "title": "RAG Group Meeting",
                "format": "pptx",
                "uri": "artifact://deck123",
                "path": str(pptx_path),
                "mime": PPTX_MIME,
                "size_bytes": pptx_path.stat().st_size,
            }
        ],
    }
    presentation = task_presentation_from_result(
        subtask_id="task-pptx",
        skill="research-pptx",
        status="succeeded",
        result=result,
    )

    cli = presentation.to_markdown()
    envelope = delivery_envelope_from_presentation(presentation)

    assert str(pptx_path) in cli
    assert "artifact://deck123" not in cli
    # Chat announces the deck once, and that line carries both the name the
    # attachment beside the message is called and the path for the owner of the
    # machine that built it. Saying the name under one heading and the path
    # under another left the reader matching the two up by title.
    message = envelope.parts[0].text
    assert f"RAG Group Meeting (pptx, 4 B): `{pptx_path}`" in message
    assert message.count(str(pptx_path)) == 1
    assert [(part.kind, part.path) for part in envelope.parts[1:]] == [("file", str(pptx_path))]
