"""Contract tests for contributed built-in skills."""

from __future__ import annotations

import asyncio
import importlib.util
import sys
import threading
from pathlib import Path
from types import SimpleNamespace

import pytest
import yaml

from omni.agent.capabilities import WORKFLOW_CAPABILITIES, deliverables_from_capabilities
from omni.agent.intent_plan import IntentType
from omni.agent.model_planner import ModelPlanProposal
from omni.agent.plan_validator import PlanValidator
from omni.agent.planner import IntentPlanner
from omni.config import load_settings
from omni.config.settings import _env_layer
from omni.skills_runtime.registry import SkillRegistry

SKILLS_ROOT = Path(__file__).resolve().parents[3] / "skills"
CONTRIBUTED_SKILLS = ("research-pptx", "livefigure", "research-ideation")


def _load_module(skill_name: str, filename: str, module_name: str):
    path = SKILLS_ROOT / skill_name / filename
    spec = importlib.util.spec_from_file_location(module_name, path)
    assert spec is not None and spec.loader is not None
    module = importlib.util.module_from_spec(spec)
    sys.modules[module_name] = module
    spec.loader.exec_module(module)
    return module


def test_new_artifact_capabilities_participate_in_workflows_and_deliverables() -> None:
    assert WORKFLOW_CAPABILITIES == {
        "literature.search",
        "artifact.figure",
        "figure.editable.pptx",
        "slides.generate",
        "synthesis.final",
        "draft.section",
        "draft.manuscript",
        "research.ideation",
        "paper.fetch.arxiv",
        "review.paper",
        "review.response",
        "poster.scientific",
    }
    assert deliverables_from_capabilities(["figure.editable.pptx"]) == ["artifact.pptx"]
    assert deliverables_from_capabilities(["slides.generate"]) == ["artifact.slides"]
    assert deliverables_from_capabilities(["review.paper"]) == ["review"]
    assert deliverables_from_capabilities(["review.response"]) == ["response_letter"]
    assert deliverables_from_capabilities(["poster.scientific"]) == ["artifact.poster"]


@pytest.mark.parametrize(
    ("capability", "expected_skill"),
    [
        ("figure.editable.pptx", "livefigure"),
        ("slides.generate", "research-pptx"),
        ("research.ideation", "research-ideation"),
        ("review.paper", "paper-review"),
        ("review.response", "review-response"),
        ("poster.scientific", "scientific-poster"),
    ],
)
def test_model_capability_routes_to_the_specialized_provider(
    capability: str,
    expected_skill: str,
) -> None:
    registry = SkillRegistry(load_settings())
    registry.build_index()
    if capability == "figure.editable.pptx":
        registry.use_admission_services(
            {"vlm": SimpleNamespace(available=True, setup_command="omni config vlm")}
        )
    planner = IntentPlanner(registry)
    plan = planner.plan_from_proposal(
        "Create the requested research artifact about RAG factuality.",
        ModelPlanProposal(
            intent_type="single_skill_task",
            required_capabilities=[capability],
            outputs=["answer"] if capability == "research.ideation" else ["artifact"],
            confidence=0.9,
            rationale="capability routing contract",
        ),
        task_id="route-specialized-skill",
    )

    assert plan.intent_type == IntentType.SINGLE_SKILL_TASK
    assert [selection.skill for selection in plan.selected_skills] == [expected_skill]
    validation = PlanValidator(registry).validate(plan)
    assert validation.ok


def test_explicit_research_pptx_binds_the_user_instruction() -> None:
    registry = SkillRegistry(load_settings())
    registry.build_index()
    plan = IntentPlanner(registry).boundary_plan(
        "$research-pptx Create a 12-slide RAG group-meeting deck.",
        task_id="explicit-research-pptx",
    )

    assert plan is not None
    validation = PlanValidator(registry).validate(plan)
    assert validation.ok
    params = plan.provider_inputs["research-pptx"]
    assert "RAG group-meeting deck" in str(params.get("topic") or params.get("input") or "")
    assert "artifact.slides" in plan.verification_plan.required_outputs


@pytest.mark.parametrize(
    ("skill_name", "instruction_field", "user_request"),
    [
        ("scientific-poster", "input", "Create a conference poster from paper.pdf."),
        ("paper-review", "input", "Review paper.pdf for NeurIPS."),
        ("review-response", "input", "Draft a response to the reviewer comments."),
    ],
)
def test_explicit_new_skill_binds_the_user_instruction(
    skill_name: str,
    instruction_field: str,
    user_request: str,
) -> None:
    registry = SkillRegistry(load_settings())
    registry.build_index()
    plan = IntentPlanner(registry).boundary_plan(
        f"${skill_name} {user_request}",
        task_id=f"explicit-{skill_name}",
    )

    assert plan is not None
    validation = PlanValidator(registry).validate(plan)
    assert validation.ok
    assert user_request in str(plan.provider_inputs[skill_name][instruction_field])


def test_scientific_poster_omni_workspace_ignores_external_output_dir(tmp_path) -> None:
    module = _load_module(
        "scientific-poster",
        "engine.py",
        "contributed_scientific_poster_engine",
    )

    class Paths:
        artifacts_dir = tmp_path / "artifacts"

    class Context:
        paths = Paths()
        session_id = "session"
        task_id = "task"

    workspace = module._runtime_io.create_workspace(  # noqa: SLF001
        {"output_dir": str(tmp_path / "outside")},
        Context(),
    )

    assert workspace.is_relative_to(Paths.artifacts_dir)


def test_contributed_manifests_declare_omni_instruction_contracts() -> None:
    manifests = {
        name: yaml.safe_load((SKILLS_ROOT / name / "SKILL.md").read_text(encoding="utf-8").split("---", 2)[1])
        for name in CONTRIBUTED_SKILLS
    }
    livefigure = manifests["livefigure"]["metadata"]["helixforge"]
    assert livefigure["input_schema"]["properties"]["input"]["x-omni"]["semantic_role"] == "instruction"
    assert livefigure["workflow"]["failure_policy"] == "continue_with_partial"
    assert "artifact.figure" in livefigure["capabilities"]
    assert livefigure["deliverables"] == ["artifact.figure", "artifact.pptx"]

    ideation = manifests["research-ideation"]["metadata"]["helixforge"]
    assert ideation["input_schema"]["properties"]["input"]["x-omni"]["semantic_role"] == "instruction"

    presentation = manifests["research-pptx"]["metadata"]["helixforge"]
    assert "artifact.slides" in presentation["capabilities"]
    assert presentation["input_schema"]["properties"]["topic"]["x-omni"]["semantic_role"] == "instruction"


def test_research_pptx_rejects_removed_export_action() -> None:
    module = _load_module("research-pptx", "engine.py", "contributed_research_pptx")
    result = module.ResearchPptxEngine.validate_params(
        arguments={"action": "export", "export_pptx_uri": "deck.pptx"}
    )
    assert result is not None
    assert result["error_info"]["code"] == "unsupported_action"


def test_livefigure_rejects_dynamic_import_escape() -> None:
    skill_dir = SKILLS_ROOT / "livefigure"
    sys.path.insert(0, str(skill_dir))
    try:
        from livefigure.pipeline import LiveFigureError, _validate_code

        code = 'getattr(__builtins__, "__import__")("os").system("id")\n'
        with pytest.raises(LiveFigureError, match="forbidden"):
            _validate_code(code)
    finally:
        sys.path.remove(str(skill_dir))


def test_livefigure_rejects_indirect_module_graph_escape() -> None:
    skill_dir = SKILLS_ROOT / "livefigure"
    sys.path.insert(0, str(skill_dir))
    try:
        from livefigure.pipeline import LiveFigureError, _validate_code

        code = (
            "import json\n"
            "json.codecs.sys.modules['os'].remove('outside')\n"
        )
        with pytest.raises(LiveFigureError, match="forbidden"):
            _validate_code(code)
    finally:
        sys.path.remove(str(skill_dir))


@pytest.mark.parametrize(
    "code",
    [
        "from pptx.opc.serialized import os as x\nx.link('a', 'b')\n",
        "from pptx.compat import sys as x\nx.modules['os'].symlink('a', 'b')\n",
        "import pptx.opc.serialized as x\nx.os.remove('outside')\n",
    ],
)
def test_livefigure_rejects_internal_pptx_import_escape(code: str) -> None:
    skill_dir = SKILLS_ROOT / "livefigure"
    sys.path.insert(0, str(skill_dir))
    try:
        from livefigure.pipeline import LiveFigureError, _validate_code

        with pytest.raises(LiveFigureError, match="forbidden"):
            _validate_code(code)
    finally:
        sys.path.remove(str(skill_dir))


@pytest.mark.asyncio
async def test_livefigure_execute_code_applies_sandbox_prefix(
    monkeypatch: pytest.MonkeyPatch, tmp_path: Path
) -> None:
    """The child runs isolated and, when supplied, under the OS sandbox prefix.

    Confinement replaces the old in-process audit hook: the host adapter passes
    a real seatbelt/bwrap prefix and the core prepends it to the child argv
    (plus ``-I`` isolation). Deterministic and env-independent — the spawn is
    intercepted, so no real backend is required.
    """
    skill_dir = SKILLS_ROOT / "livefigure"
    sys.path.insert(0, str(skill_dir))
    try:
        import asyncio as _asyncio

        from livefigure.pipeline import _execute_code

        script = tmp_path / "livefigure.py"
        script.write_text("pass\n", encoding="utf-8")
        pptx_path = tmp_path / "livefigure.pptx"
        prefix = ("sandbox-exec", "-p", "(version 1)(deny file-write*)")
        captured: dict[str, tuple[str, ...]] = {}

        class _FakeProc:
            returncode = 0

            async def communicate(self) -> tuple[bytes, bytes]:
                pptx_path.write_bytes(b"PK\x03\x04")
                return (b"", b"")

        async def _fake_exec(*argv: str, **_kwargs: object) -> _FakeProc:
            captured["argv"] = argv
            return _FakeProc()

        monkeypatch.setattr(_asyncio, "create_subprocess_exec", _fake_exec)
        await _execute_code(script, tmp_path, pptx_path, sandbox_prefix=prefix)

        argv = captured["argv"]
        assert argv[: len(prefix)] == prefix
        assert "-I" in argv
        boot = Path(str(argv[-1]))
        assert boot.name == "_omni_livefigure_boot.py"
        boot_src = boot.read_text(encoding="utf-8")
        assert "runpy.run_path" in boot_src
        assert repr(str(script.resolve())) in boot_src
    finally:
        sys.path.remove(str(skill_dir))


def test_livefigure_rejects_writes_outside_output_file() -> None:
    skill_dir = SKILLS_ROOT / "livefigure"
    sys.path.insert(0, str(skill_dir))
    try:
        from livefigure.pipeline import LiveFigureError, _validate_code

        code = (
            "from pptx import Presentation\n"
            "presentation = Presentation()\n"
            'presentation.save("/tmp/outside.pptx")\n'
        )
        with pytest.raises(LiveFigureError, match="livefigure.pptx"):
            _validate_code(code)
    finally:
        sys.path.remove(str(skill_dir))


def test_livefigure_rejects_network_imports() -> None:
    """Network access is denied statically: ``socket`` is off the import allowlist.

    The old runtime relied on an audit hook to trap ``socket()`` at call time;
    the code denylist now rejects the import outright, so hostile code never
    reaches execution regardless of OS-sandbox availability.
    """
    skill_dir = SKILLS_ROOT / "livefigure"
    sys.path.insert(0, str(skill_dir))
    try:
        from livefigure.pipeline import LiveFigureError, _validate_code

        with pytest.raises(LiveFigureError, match="forbidden"):
            _validate_code("import socket\nsocket.socket()\n")
    finally:
        sys.path.remove(str(skill_dir))


def test_livefigure_allows_bundled_tools_helper() -> None:
    """The prompt's ``from tools import *`` (and documented helpers) validates.

    LiveFigure copies a trusted ``tools.py`` beside the generated script and
    tells the model to star-import it, so a realistic drawing script that uses
    the helpers and saves ``livefigure.pptx`` must pass the code denylist.
    """
    skill_dir = SKILLS_ROOT / "livefigure"
    sys.path.insert(0, str(skill_dir))
    try:
        from livefigure.pipeline import _validate_code

        code = (
            "from pptx import Presentation\n"
            "from tools import *\n"
            "presentation = Presentation()\n"
            "slide = presentation.slides.add_slide(presentation.slide_layouts[6])\n"
            "a = add_block(slide, 1, 1, 2, 1, text='A')\n"
            "b = add_block(slide, 4, 1, 2, 1, text='B')\n"
            "add_connector(slide, a, b)\n"
            "presentation.save('livefigure.pptx')\n"
        )
        _validate_code(code)  # must not raise
    finally:
        sys.path.remove(str(skill_dir))


def test_livefigure_tools_helper_import_is_scoped() -> None:
    """The ``tools`` allowance is narrow: star/helpers only, no reflection.

    ``import tools`` (root) and importing the module's re-exported internals
    (``os``, ``requests`` …) must stay forbidden so the trusted helper cannot
    become an escape hatch for network/filesystem access.
    """
    skill_dir = SKILLS_ROOT / "livefigure"
    sys.path.insert(0, str(skill_dir))
    try:
        from livefigure.pipeline import LiveFigureError, _validate_code

        with pytest.raises(LiveFigureError, match="forbidden"):
            _validate_code("import tools\ntools.add_block()\n")
        with pytest.raises(LiveFigureError, match="forbidden"):
            _validate_code("from tools import os\nos.system('id')\n")
        with pytest.raises(LiveFigureError, match="forbidden"):
            _validate_code("from tools import requests\nrequests.get('http://x')\n")
    finally:
        sys.path.remove(str(skill_dir))


def test_livefigure_bundled_tools_exports_only_drawing_helpers() -> None:
    """``tools.__all__`` must not re-export network/OS names via ``import *``.

    A missing/loosened ``__all__`` would let ``from tools import *`` leak
    ``os``/``requests``/``Image`` into generated code, past the denylist. Parsed
    from source (no import) so it never depends on the helper's own third-party
    deps being installed in the test environment.
    """
    import ast

    source = (SKILLS_ROOT / "livefigure" / "livefigure" / "tools.py").read_text(encoding="utf-8")
    exported: set[str] | None = None
    for node in ast.walk(ast.parse(source)):
        if isinstance(node, ast.Assign) and any(
            isinstance(t, ast.Name) and t.id == "__all__" for t in node.targets
        ):
            exported = {
                elt.value
                for elt in node.value.elts  # type: ignore[attr-defined]
                if isinstance(elt, ast.Constant)
            }
    assert exported == {
        "add_block",
        "add_connector",
        "add_container",
        "add_custom_route_arrow",
        "add_free_arrow",
        "add_label",
    }
    assert {"os", "requests", "Image", "json", "base64", "io"}.isdisjoint(exported)


def test_livefigure_vendored_security_contract_is_intact() -> None:
    """Canary against a vendor re-sync silently regressing the code executor.

    LiveFigure's generated-code path is only safe because of two invariants:
    (1) a static import/call denylist gates the source before execution, and
    (2) an OS ``sandbox_prefix`` seam lets the host add kernel write-confinement.
    The abandoned in-process audit hook (``sandbox_runner.py``) must stay gone.
    A future ``scy`` merge that drops any of these should fail loudly here
    rather than ship an ungated executor.
    """
    import dataclasses

    skill_dir = SKILLS_ROOT / "livefigure"
    assert not (skill_dir / "livefigure" / "sandbox_runner.py").exists()
    sys.path.insert(0, str(skill_dir))
    try:
        from livefigure.pipeline import (
            _ALLOWED_IMPORT_ROOTS,
            _BLOCKED_CALLS,
            LiveFigureError,
            PipelineConfig,
            _validate_code,
        )

        assert "pptx" in _ALLOWED_IMPORT_ROOTS and "os" not in _ALLOWED_IMPORT_ROOTS
        assert {"eval", "exec", "getattr", "open"} <= _BLOCKED_CALLS
        with pytest.raises(LiveFigureError, match="forbidden"):
            _validate_code("import os\nos.system('id')\n")

        assert "sandbox_prefix" in {f.name for f in dataclasses.fields(PipelineConfig)}
    finally:
        sys.path.remove(str(skill_dir))


def test_unpublished_livefigure_gemini_environment_is_ignored(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    monkeypatch.setenv("OMNI_LIVEFIGURE_GEMINI_BASE_URL", "https://example.invalid/v1beta")
    monkeypatch.setenv("OMNI_LIVEFIGURE_GEMINI_API_KEY", "secret-value")
    layer = _env_layer()
    assert "livefigure" not in layer


def test_livefigure_omni_engine_ignores_untrusted_output_dir(tmp_path: Path) -> None:
    module = _load_module("livefigure", "engine.py", "contributed_livefigure_engine")

    class Paths:
        artifacts_dir = tmp_path / "artifacts"

    class Context:
        scratch_dir = tmp_path / "exec"
        output_dir = tmp_path / "outbox"
        paths = Paths()

    resolved = module._resolve_output_dir(Context(), {"output_dir": "/tmp/outside"})
    assert resolved.is_relative_to(Context.scratch_dir)
    assert resolved.as_posix().endswith(f"livefigure-runs/{resolved.name}")
    assert not resolved.is_relative_to(Paths.artifacts_dir)
    assert "/tmp/outside" not in resolved.as_posix()


def test_livefigure_omni_engine_does_not_use_artifacts_dir(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    """Control-store artifacts_dir is never a generated-code cwd (sandbox would deny it)."""
    module = _load_module("livefigure", "engine.py", "contributed_livefigure_engine")
    monkeypatch.chdir(tmp_path)
    monkeypatch.setenv("TMPDIR", str(tmp_path / "tmp"))

    class Paths:
        artifacts_dir = tmp_path / "artifacts"

    class Context:
        paths = Paths()

    resolved = module._resolve_output_dir(Context(), {})
    assert not resolved.is_relative_to(Paths.artifacts_dir)
    assert resolved.is_relative_to(tmp_path / "tmp")


@pytest.mark.asyncio
async def test_research_ideation_runs_blocking_pipeline_off_event_loop(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    module = _load_module("research-ideation", "engine.py", "contributed_research_ideation")
    main_thread = threading.get_ident()

    def fake_pipeline(**_kwargs):
        return {
            "status": "ok",
            "summary": "done",
            "worker_thread": threading.get_ident(),
            "steps": {},
        }

    monkeypatch.setattr(module._core, "run_pipeline", fake_pipeline)
    monkeypatch.setattr(
        module,
        "_resolve_s2_key",
        lambda _ctx: "scoped-s2-secret",
    )
    engine = module.ResearchIdeationEngine()
    engine.ctx = SimpleNamespace(llm=object(), artifacts=None, paths=None)
    result = await asyncio.wait_for(
        engine.execute(input="RAG factuality", use_tools=False),
        timeout=2,
    )
    assert result["worker_thread"] != main_thread
