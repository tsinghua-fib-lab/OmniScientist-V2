import base64
import io
import json
import os

import requests
from PIL import Image
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls
from pptx.util import Inches, Pt

GEMINI_API_URL = os.getenv("LIVEFIGURE_GEMINI_IMAGE_URL", "")
GEMINI_API_KEY = os.getenv("LIVEFIGURE_GEMINI_API_KEY", "")

# Public drawing helpers only. Generated code runs ``from tools import *``; a
# narrow ``__all__`` keeps that star import from re-exporting the module's
# imported names (``os``, ``requests``, ``Image`` …) into the generated script,
# so the code denylist stays effective (no network/OS egress via this helper).
__all__ = [
    "add_block",
    "add_connector",
    "add_container",
    "add_custom_route_arrow",
    "add_free_arrow",
    "add_label",
]


def _apply_arrow_xml(connector, end_arrow=True, size="med"):
    ln = connector.line._get_or_add_ln()
    arrow_xml = (
        f'<a:tailEnd xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
        f'type="triangle" w="{size}" len="{size}"/>'
    )
    if not end_arrow:
        arrow_xml = arrow_xml.replace("tailEnd", "headEnd")
    ln.append(parse_xml(arrow_xml))


def _apply_gradient_xml(connector, color_start, color_end):
    ln = connector.line._get_or_add_ln()
    for child in list(ln):
        if "Fill" in child.tag:
            ln.remove(child)
    gradient_xml = f"""
    <a:gradFill {nsdecls('a')} flip="0" rotWithShape="1">
      <a:gsLst>
        <a:gs pos="0"><a:srgbClr val="{color_start}"/></a:gs>
        <a:gs pos="100000"><a:srgbClr val="{color_end}"/></a:gs>
      </a:gsLst>
      <a:lin ang="0" scaled="1"/>
    </a:gradFill>
    """
    ln.insert(0, parse_xml(gradient_xml))


def _parse_color(color_input):
    if isinstance(color_input, str):
        hex_color = color_input.lstrip("#")
        if not hex_color:
            return RGBColor(0, 0, 0)
        return RGBColor(int(hex_color[:2], 16), int(hex_color[2:4], 16), int(hex_color[4:], 16))
    if isinstance(color_input, RGBColor):
        return color_input
    return RGBColor(0, 0, 0)


def _set_shape_alpha(shape, alpha):
    if alpha >= 1.0:
        return
    fill = shape.fill
    if not hasattr(fill, "fore_color"):
        return
    try:
        xfill = fill.fore_color._xFill
        srgb_clr = xfill.srgbClr
        if srgb_clr is not None:
            alpha_tag = "{http://schemas.openxmlformats.org/drawingml/2006/main}alpha"
            for elem in srgb_clr.findall(alpha_tag):
                srgb_clr.remove(elem)
            alpha_int = int(alpha * 100000)
            alpha_elem = parse_xml(f'<a:alpha {nsdecls("a")} val="{alpha_int}"/>')
            srgb_clr.append(alpha_elem)
    except Exception:
        return


def _call_gemini_strict(prompt, aspect_ratio="1:1"):
    refined_prompt = (
        f"{prompt}. Vector icon style, flat design, minimal, "
        "solid white background, isolated, no shadows, clean edges."
    )
    payload = json.dumps(
        {
            "contents": [{"role": "user", "parts": [{"text": refined_prompt}]}],
            "generationConfig": {
                "responseModalities": ["IMAGE"],
                "imageConfig": {"aspectRatio": aspect_ratio},
            },
        }
    )
    headers = {"Authorization": f"Bearer {GEMINI_API_KEY}", "Content-Type": "application/json"}
    response = requests.post(GEMINI_API_URL, headers=headers, data=payload)
    response.raise_for_status()
    response_json = response.json()
    if "candidates" in response_json and response_json["candidates"]:
        for part in response_json["candidates"][0]["content"]["parts"]:
            if "inlineData" in part:
                return io.BytesIO(base64.b64decode(part["inlineData"]["data"]))
    return None


def _remove_white_background(image_stream, tolerance=50, crop_tight=True):
    img = Image.open(image_stream).convert("RGBA")
    datas = img.getdata()
    new_data = []
    for item in datas:
        if item[0] > (255 - tolerance) and item[1] > (255 - tolerance) and item[2] > (255 - tolerance):
            new_data.append((255, 255, 255, 0))
        else:
            new_data.append(item)
    img.putdata(new_data)
    if crop_tight:
        bbox = img.getbbox()
        if bbox:
            img = img.crop(bbox)
    output_stream = io.BytesIO()
    img.save(output_stream, format="PNG")
    output_stream.seek(0)
    return output_stream


def add_connector(
    slide,
    source_shape,
    dest_shape,
    type="curve",
    color="1F4E79",
    width=3.0,
    gradient_start=None,
    gradient_end=None,
    arrow_size=None,
    conn_src=None,
    conn_dest=None,
):
    s_x = source_shape.left + (source_shape.width // 2)
    s_y = source_shape.top + (source_shape.height // 2)
    d_x = dest_shape.left + (dest_shape.width // 2)
    d_y = dest_shape.top + (dest_shape.height // 2)

    if conn_src is not None and conn_dest is not None:
        src_idx = int(conn_src)
        dst_idx = int(conn_dest)
    else:
        dx, dy = d_x - s_x, d_y - s_y
        if abs(dx) > abs(dy):
            src_idx, dst_idx = (3, 1) if dx > 0 else (1, 3)
        else:
            src_idx, dst_idx = (2, 0) if dy > 0 else (0, 2)

    type_map = {"curve": MSO_CONNECTOR.CURVE, "elbow": MSO_CONNECTOR.ELBOW, "straight": MSO_CONNECTOR.STRAIGHT}
    connector = slide.shapes.add_connector(type_map.get(type, MSO_CONNECTOR.CURVE), s_x, s_y, d_x, d_y)
    connector.begin_connect(source_shape, src_idx)
    connector.end_connect(dest_shape, dst_idx)

    connector.line.width = Pt(width)
    if gradient_start and gradient_end:
        _apply_gradient_xml(connector, gradient_start.lstrip("#"), gradient_end.lstrip("#"))
    else:
        connector.line.color.rgb = _parse_color(color)

    if arrow_size is None:
        arrow_size = "sm" if width < 2.0 else "med" if width <= 4.5 else "lg"
    _apply_arrow_xml(connector, end_arrow=True, size=arrow_size)
    return connector


def add_free_arrow(
    slide,
    start_x,
    start_y,
    end_x,
    end_y,
    type="straight",
    color="1F4E79",
    width=3.0,
    gradient_start=None,
    gradient_end=None,
):
    type_map = {"curve": MSO_CONNECTOR.CURVE, "elbow": MSO_CONNECTOR.ELBOW, "straight": MSO_CONNECTOR.STRAIGHT}
    connector = slide.shapes.add_connector(
        type_map.get(type, MSO_CONNECTOR.STRAIGHT),
        Inches(start_x),
        Inches(start_y),
        Inches(end_x),
        Inches(end_y),
    )
    connector.line.width = Pt(width)
    if gradient_start and gradient_end:
        _apply_gradient_xml(connector, gradient_start.lstrip("#"), gradient_end.lstrip("#"))
    else:
        connector.line.color.rgb = _parse_color(color)
    _apply_arrow_xml(connector, end_arrow=True, size="lg")
    return connector


def add_custom_route_arrow(slide, points, color="333333", width=2.0, end_arrow=True):
    if not points or len(points) < 2:
        return None
    start_x, start_y = points[0]
    builder = slide.shapes.build_freeform(Inches(start_x), Inches(start_y))
    vertex_pairs = [(Inches(x), Inches(y)) for x, y in points[1:]]
    builder.add_line_segments(vertex_pairs, close=False)
    shape = builder.convert_to_shape()
    shape.fill.background()
    shape.line.width = Pt(width)
    shape.line.color.rgb = _parse_color(color)
    if end_arrow:
        _apply_arrow_xml(shape, end_arrow=True, size="lg")
    return shape


def add_block(
    slide,
    x,
    y,
    w,
    h,
    text=None,
    fill_color="FFFFFF",
    stroke_color="000000",
    shape_type=MSO_SHAPE.ROUNDED_RECTANGLE,
    font_size=12,
    font_color="000000",
    bold=False,
    alpha=1.0,
):
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = _parse_color(fill_color)
        if alpha < 1.0:
            _set_shape_alpha(shape, alpha)
    else:
        shape.fill.background()

    if stroke_color:
        shape.line.color.rgb = _parse_color(stroke_color)
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()

    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
        p.text = str(text)
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(font_size)
        p.font.color.rgb = _parse_color(font_color)
        p.font.bold = bold
    return shape


def add_label(slide, text, x, y, w=None, h=None, font_size=12, color="000000", bold=False, align="center"):
    width_in = Inches(w) if w else Inches(2.0)
    height_in = Inches(h) if h else Inches(0.5)
    textbox = slide.shapes.add_textbox(Inches(x), Inches(y), width_in, height_in)
    tf = textbox.text_frame
    tf.word_wrap = True
    if w is None:
        tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    p = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
    p.text = str(text)
    p.font.size = Pt(font_size)
    p.font.color.rgb = _parse_color(color)
    p.font.bold = bold
    align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT, "justify": PP_ALIGN.JUSTIFY}
    p.alignment = align_map.get(str(align).lower(), PP_ALIGN.CENTER)
    return textbox


def add_container(slide, x, y, w, h, title=None, fill_color="F5F5F5", stroke_color="CCCCCC", alpha=1.0):
    container = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill_color:
        container.fill.solid()
        container.fill.fore_color.rgb = _parse_color(fill_color)
        if alpha < 1.0:
            _set_shape_alpha(container, alpha)
    else:
        container.fill.background()

    if stroke_color:
        container.line.color.rgb = _parse_color(stroke_color)
        container.line.width = Pt(1.5)
    else:
        container.line.fill.background()
    container.shadow.inherit = False

    if title:
        add_label(slide, title, x=x, y=y + 0.1, w=w, h=0.4, font_size=11, bold=True, color="333333", align="center")
    return container
