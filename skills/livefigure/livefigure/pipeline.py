"""One-pass VLM-to-PPTX pipeline without image rendering or visual critique."""

from __future__ import annotations

import ast
import asyncio
import os
import re
import shutil
import sys
import zipfile
from collections.abc import Awaitable, Callable
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Protocol

from .prompts import build_generation_prompt, build_repair_prompt
from .vlm import VlmError, reference_as_data_url
from .vlm import reference_bytes as decode_reference_bytes

# ``(stage, pct)`` at minimum; a host that speaks the stage contract also takes
# the ``stage_id``/``milestone``/``stats`` keywords, so the shape is negotiated
# per call rather than fixed here.
Progress = Callable[..., Awaitable[None]]

# --- Static code denylist -------------------------------------------------
#
# Model-generated ``python-pptx`` source is confined to a tiny allowlist of
# imports and blocked from the reflection / filesystem / process primitives
# that would let it escape (import graph walking, ``getattr(__builtins__, …)``,
# ``open``, ``os``/``subprocess``/``socket`` …). This static gate is the first
# line of defence; ``_execute_code`` adds OS-level write confinement on top.
_ALLOWED_IMPORT_ROOTS = frozenset({"pptx", "math"})
_ALLOWED_FROM_IMPORTS: dict[str, frozenset[str]] = {
    "math": frozenset(
        {
            "ceil",
            "cos",
            "degrees",
            "e",
            "exp",
            "floor",
            "log",
            "pi",
            "radians",
            "sin",
            "sqrt",
            "tan",
            "tau",
        }
    ),
    "pptx": frozenset({"Presentation"}),
    "pptx.chart.data": frozenset(
        {"BubbleChartData", "CategoryChartData", "ChartData", "XyChartData"}
    ),
    "pptx.dml.color": frozenset({"RGBColor"}),
    "pptx.enum.chart": frozenset(
        {"XL_CHART_TYPE", "XL_LABEL_POSITION", "XL_LEGEND_POSITION"}
    ),
    "pptx.enum.dml": frozenset({"MSO_LINE_DASH_STYLE", "MSO_THEME_COLOR"}),
    "pptx.enum.lang": frozenset({"MSO_LANGUAGE_ID"}),
    "pptx.enum.shapes": frozenset(
        {
            "MSO_AUTO_SHAPE_TYPE",
            "MSO_CONNECTOR",
            "MSO_CONNECTOR_TYPE",
            "MSO_SHAPE",
            "MSO_SHAPE_TYPE",
            "PP_PLACEHOLDER",
        }
    ),
    "pptx.enum.text": frozenset(
        {
            "MSO_ANCHOR",
            "MSO_AUTO_SIZE",
            "MSO_TEXT_DIRECTION",
            "MSO_VERTICAL_ANCHOR",
            "PP_ALIGN",
        }
    ),
    "pptx.util": frozenset({"Cm", "Emu", "Inches", "Pt"}),
    # The skill copies a trusted ``tools.py`` (drawing helpers) beside the
    # generated script and the prompt instructs the model to
    # ``from tools import *``. Allow the star the prompt asks for and the
    # documented helper names; ``import tools`` (root) stays off
    # ``_ALLOWED_IMPORT_ROOTS`` so generated code cannot reach the module's
    # non-exported internals reflectively (e.g. ``tools._call_gemini_strict``).
    "tools": frozenset(
        {
            "*",
            "add_block",
            "add_connector",
            "add_container",
            "add_custom_route_arrow",
            "add_free_arrow",
            "add_label",
        }
    ),
}
_BLOCKED_CALLS = frozenset(
    {
        "__import__",
        "breakpoint",
        "compile",
        "delattr",
        "dir",
        "eval",
        "exec",
        "getattr",
        "globals",
        "help",
        "input",
        "locals",
        "open",
        "setattr",
        "vars",
    }
)
_BLOCKED_ATTRIBUTES = frozenset(
    {
        "builtins",
        "codecs",
        "compile",
        "ctypes",
        "environ",
        "eval",
        "exec",
        "fork",
        "import_module",
        "modules",
        "open",
        "os",
        "popen",
        "remove",
        "rename",
        "replace",
        "rmdir",
        "shutil",
        "socket",
        "spawn",
        "subprocess",
        "sys",
        "system",
        "unlink",
    }
)


class LiveFigureError(RuntimeError):
    """A generation error that can be reported as a structured skill result."""

    def __init__(
        self,
        message: str,
        *,
        code: str = "livefigure_failed",
        category: str = "generation",
        retryable: bool = True,
    ) -> None:
        super().__init__(message)
        self.code = code
        self.category = category
        self.retryable = retryable


class VlmTextGenerator(Protocol):
    """Portable port implemented by Omni's host service or the env runner."""

    async def generate_text(
        self,
        prompt: str,
        *,
        reference_image_uri: str | None = None,
    ) -> str: ...


@dataclass(frozen=True)
class PipelineConfig:
    vlm: VlmTextGenerator | None = None
    max_code_retries: int = 1
    reference_roots: tuple[Path, ...] = field(default_factory=tuple)
    reference_files: tuple[Path, ...] = field(default_factory=tuple)
    # OS-level confinement argv prefix (e.g. seatbelt / bwrap) prepended to the
    # child that runs generated python-pptx source. Empty = run unwrapped; the
    # host adapter supplies it, keeping this core portable.
    sandbox_prefix: tuple[str, ...] = field(default_factory=tuple)


@dataclass(frozen=True)
class PipelineResult:
    title: str
    pptx_path: Path
    code_path: Path
    input_path: Path
    reference_path: Path | None
    attempts: int


async def generate_pptx(
    requirement: str,
    *,
    title: str,
    output_dir: Path,
    config: PipelineConfig,
    reference_image_uri: str | None = None,
    progress: Progress | None = None,
) -> PipelineResult:
    """Generate a PPTX, using an optional local/data-URL visual reference."""
    output_dir.mkdir(parents=True, exist_ok=True)
    shutil.copyfile(Path(__file__).with_name("tools.py"), output_dir / "tools.py")
    reference_path: Path | None = None
    vlm_client = config.vlm
    if vlm_client is None:
        raise LiveFigureError(
            "VLM configuration is not available",
            code="vlm_not_configured",
            category="configuration",
            retryable=False,
        )
    try:
        reference_data_url: str | None = None
        if reference_image_uri:
            reference_data_url = reference_as_data_url(
                reference_image_uri,
                allowed_roots=config.reference_roots,
                allowed_files=config.reference_files,
            )
            raw, reference_mime = decode_reference_bytes(reference_data_url)
            suffix = _reference_suffix(reference_mime)
            reference_path = output_dir / f"reference{suffix}"
            reference_path.write_bytes(raw)

        await _progress(progress, "generate pptx code", 0.40)
        code = await vlm_client.generate_text(
            _code_prompt(requirement, title),
            reference_image_uri=reference_data_url,
        )
    except VlmError as exc:
        raise _from_vlm_error(exc) from exc

    input_path = output_dir / "input.txt"
    input_path.write_text(requirement, encoding="utf-8")
    max_attempts = max(1, int(config.max_code_retries) + 1)
    last_error = LiveFigureError("PPTX generation failed")
    for attempt in range(1, max_attempts + 1):
        await _progress(progress, "validate and build pptx", 0.45 + (attempt / max_attempts) * 0.45)
        code = _strip_fence(code)
        code_path = output_dir / "generated_figure.py"
        code_path.write_text(code, encoding="utf-8")
        pptx_path = output_dir / "livefigure.pptx"
        try:
            _validate_code(code)
            await _execute_code(
                code_path, output_dir, pptx_path, sandbox_prefix=config.sandbox_prefix
            )
            _validate_pptx(pptx_path)
            await _progress(
                progress,
                "pptx ready",
                1.0,
                stage_id="livefigure.done",
                milestone="Live figure built",
                stats={"attempts": attempt},
            )
            return PipelineResult(title, pptx_path, code_path, input_path, reference_path, attempt)
        except LiveFigureError as exc:
            last_error = exc
            if attempt == max_attempts:
                break
            try:
                repair_prompt = _repair_prompt(requirement, title, code, str(last_error))
                code = await vlm_client.generate_text(repair_prompt)
            except VlmError as repair_exc:
                last_error = _from_vlm_error(repair_exc)
                break
    raise LiveFigureError(
        f"PPTX generation failed: {last_error}",
        code=last_error.code,
        category=last_error.category,
        retryable=last_error.retryable,
    )


def _from_vlm_error(exc: VlmError) -> LiveFigureError:
    return LiveFigureError(
        str(exc),
        code=exc.code,
        category=exc.category,
        retryable=exc.retryable,
    )


def _reference_suffix(mime: str) -> str:
    return {
        "image/jpeg": ".jpg",
        "image/gif": ".gif",
        "image/webp": ".webp",
    }.get(mime.lower(), ".png")


def _code_prompt(requirement: str, title: str) -> str:
    return build_generation_prompt(requirement, title)


def _repair_prompt(requirement: str, title: str, code: str, error: str) -> str:
    return build_repair_prompt(requirement, title, code, error)


def _strip_fence(code: str) -> str:
    text = code.strip()
    fenced = re.search(r"```(?:python)?\s*\n?(.*?)```", text, flags=re.DOTALL | re.IGNORECASE)
    if fenced:
        text = fenced.group(1)
    return text.strip() + "\n"


def _validate_code(code: str) -> None:
    if not code or len(code) > 60_000:
        raise LiveFigureError("Generated PPTX code is empty or too large")
    try:
        tree = ast.parse(code)
    except SyntaxError as exc:
        raise LiveFigureError(f"Generated PPTX code has invalid syntax: {exc.msg}") from exc
    for node in ast.walk(tree):
        if isinstance(node, ast.Import):
            if any(alias.name not in _ALLOWED_IMPORT_ROOTS for alias in node.names):
                raise LiveFigureError("Generated code contains a forbidden import")
        elif isinstance(node, ast.ImportFrom):
            module = node.module or ""
            allowed_names = _ALLOWED_FROM_IMPORTS.get(module)
            imported_names = {alias.name for alias in node.names}
            if allowed_names is None or not imported_names <= allowed_names:
                raise LiveFigureError("Generated code contains a forbidden import")
        elif (
            isinstance(node, ast.Call)
            and isinstance(node.func, ast.Name)
            and node.func.id in _BLOCKED_CALLS
        ):
            raise LiveFigureError(f"Generated code contains a forbidden call: {node.func.id}")
        elif (
            isinstance(node, ast.Call)
            and isinstance(node.func, ast.Attribute)
            and node.func.attr == "save"
        ):
            if (
                len(node.args) != 1
                or not isinstance(node.args[0], ast.Constant)
                or node.args[0].value != "livefigure.pptx"
                or node.keywords
            ):
                raise LiveFigureError("Generated code must save only to livefigure.pptx")
        elif isinstance(node, ast.Attribute) and node.attr in {"head_end", "tail_end"}:
            raise LiveFigureError(
                "Generated code uses arrowhead APIs unsupported by this python-pptx version"
            )
        elif isinstance(node, ast.Attribute) and node.attr.lower() in _BLOCKED_ATTRIBUTES:
            raise LiveFigureError(
                f"Generated code contains a forbidden attribute: {node.attr}"
            )
        elif isinstance(node, ast.Name) and node.id == "MSO_LINE":
            raise LiveFigureError("Generated code uses the unsupported MSO_LINE arrow API")
        elif isinstance(node, ast.Attribute) and node.attr.startswith("__"):
            raise LiveFigureError("Generated code contains a forbidden dunder attribute")
        elif isinstance(node, ast.Name) and node.id.startswith("__"):
            raise LiveFigureError("Generated code contains a forbidden dunder name")


def _isolated_script_command(code_path: Path, output_dir: Path) -> list[str]:
    """Run generated code isolated, but keep the sibling helper importable."""
    boot_path = output_dir / "_omni_livefigure_boot.py"
    boot_path.write_text(
        "import runpy, sys\n"
        f"sys.path.insert(0, {str(output_dir.resolve())!r})\n"
        f"runpy.run_path({str(code_path.resolve())!r}, run_name='__main__')\n",
        encoding="utf-8",
    )
    return [sys.executable, "-I", str(boot_path)]


async def _execute_code(
    code_path: Path,
    output_dir: Path,
    pptx_path: Path,
    *,
    sandbox_prefix: tuple[str, ...] = (),
) -> None:
    if pptx_path.exists():
        pptx_path.unlink()
    env = {"PATH": os.environ.get("PATH", ""), "PYTHONNOUSERSITE": "1"}
    # Isolated interpreter (``-I``): ignore env vars / user site / cwd on
    # sys.path. That also hides the trusted ``tools.py`` copied beside the
    # script, so a host-owned boot fragment reinserts only ``output_dir``.
    # When the host supplies a ``sandbox_prefix`` (seatbelt / bwrap), the
    # child additionally runs under kernel write-confinement.
    argv = [*sandbox_prefix, *_isolated_script_command(code_path, output_dir)]
    try:
        proc = await asyncio.create_subprocess_exec(
            *argv,
            cwd=str(output_dir),
            env=env,
            stdout=asyncio.subprocess.PIPE,
            stderr=asyncio.subprocess.PIPE,
        )
        _, stderr = await asyncio.wait_for(proc.communicate(), timeout=75)
    except TimeoutError:
        if proc.returncode is None:
            proc.kill()
            await proc.wait()
        raise LiveFigureError("PPTX code execution timed out after 75 seconds") from None
    except OSError as exc:
        raise LiveFigureError(f"Could not execute PPTX code: {exc}") from exc
    if proc.returncode != 0:
        raise LiveFigureError(
            stderr.decode("utf-8", "replace")[-3000:] or "PPTX code execution failed"
        )
    if not pptx_path.is_file():
        raise LiveFigureError("Generated code did not create the required livefigure.pptx")


def _validate_pptx(path: Path) -> None:
    try:
        with zipfile.ZipFile(path) as archive:
            if (
                "[Content_Types].xml" not in archive.namelist()
                or "ppt/presentation.xml" not in archive.namelist()
            ):
                raise LiveFigureError("Generated file is not a valid PPTX")
    except zipfile.BadZipFile as exc:
        raise LiveFigureError("Generated file is not a valid PPTX") from exc
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        presentation = Presentation(str(path))
    except Exception as exc:  # noqa: BLE001 - normalize parser/dependency failures
        raise LiveFigureError("Generated file could not be parsed as a PPTX") from exc
    if len(presentation.slides) != 1:
        raise LiveFigureError("LiveFigure output must contain exactly one slide")
    shapes = list(presentation.slides[0].shapes)
    if not shapes or not any(shape.shape_type != MSO_SHAPE_TYPE.PICTURE for shape in shapes):
        raise LiveFigureError(
            "LiveFigure output must contain at least one editable PowerPoint shape"
        )


async def _progress(callback: Progress | None, stage: str, pct: float, **data: Any) -> None:
    if callback is not None:
        try:
            await callback(stage, pct, **data)
        except TypeError:
            await callback(stage, pct)
