#!/usr/bin/env python3
"""Render a canonical scientific-poster scene as native PowerPoint objects."""

from __future__ import annotations

import argparse
import base64
import json
import math
import sys
from io import BytesIO
from pathlib import Path
from urllib.parse import unquote, urlparse

from lxml import etree

SKILL_DIR = Path(__file__).resolve().parents[1]
if str(SKILL_DIR) not in sys.path:
    sys.path.insert(0, str(SKILL_DIR))

from posterlib.delivery.pptx_math import (  # noqa: E402
    A14_NAMESPACE,
    build_text_math_xml,
)
from posterlib.delivery.pptx_scene import normalize_scene  # noqa: E402

MC_NAMESPACE = "http://schemas.openxmlformats.org/markup-compatibility/2006"
_SVG_TARGET_PPI = 300
_SVG_MAX_LONG_EDGE_PIXELS = 8192
_SVG_MAX_PIXELS = 24_000_000


def render_scene(scene_path: Path, output_path: Path) -> None:
    """Write one editable custom-size slide from a validated scene."""

    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
        from pptx.util import Inches, Pt
    except ImportError as exc:
        raise RuntimeError("python-pptx is required for editable PPTX export") from exc

    scene = normalize_scene(json.loads(scene_path.read_text(encoding="utf-8")))
    presentation = Presentation()
    presentation.slide_width = Inches(scene["page"]["width_in"])
    presentation.slide_height = Inches(scene["page"]["height_in"])
    properties = presentation.core_properties
    properties.author = "OmniScientist scientific-poster"
    properties.subject = "Editable academic poster"
    properties.title = "Scientific poster"

    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = RGBColor.from_string(scene["page"]["background"])
    alignments = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
    }

    for item in scene["objects"]:
        x, y, width, height = (item[key] for key in ("x", "y", "w", "h"))
        geometry = tuple(Inches(value) for value in (x, y, width, height))
        if item["kind"] == "shape":
            shape_type = (
                MSO_SHAPE.ROUNDED_RECTANGLE
                if item["radius"] > 0
                else MSO_SHAPE.RECTANGLE
            )
            shape = slide.shapes.add_shape(shape_type, *geometry)
            if item["radius"] > 0:
                shape.adjustments[0] = min(
                    0.5,
                    float(item["radius"]) / min(float(width), float(height)),
                )
            if item["fill_enabled"]:
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor.from_string(item["fill"])
            else:
                shape.fill.background()
            if item["line_width_pt"] > 0:
                shape.line.color.rgb = RGBColor.from_string(item["line"])
                shape.line.width = Pt(item["line_width_pt"])
            else:
                shape.line.fill.background()
            _remove_theme_shape_style(shape)
        elif item["kind"] == "image":
            source = _image_source(
                item["src"],
                scene_path.parent,
                width_in=width,
                height_in=height,
            )
            shape = slide.shapes.add_picture(source, *geometry)
            if item["alt"]:
                shape._element.nvPicPr.cNvPr.set("descr", item["alt"])
        elif item["kind"] == "equation":
            _add_native_equation(
                slide,
                item,
                geometry,
                scene_path.parent,
                alignments,
                RGBColor,
                MSO_ANCHOR,
                Pt,
            )
            continue
        elif item["kind"] == "table":
            rows = item["rows"]
            shape = slide.shapes.add_table(len(rows), len(rows[0]), *geometry)
            for row, height_in in zip(
                shape.table.rows,
                item["row_heights"],
                strict=True,
            ):
                row.height = Inches(height_in)
            for column, width_in in zip(
                shape.table.columns,
                item["column_widths"],
                strict=True,
            ):
                column.width = Inches(width_in)
            for row_index, row in enumerate(rows):
                for column_index, value in enumerate(row):
                    cell = shape.table.cell(row_index, column_index)
                    cell.text = value
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor.from_string(
                        item["header_fill"] if row_index == 0 else item["body_fill"]
                    )
                    text_frame = cell.text_frame
                    text_frame.margin_left = Inches(0.04)
                    text_frame.margin_right = Inches(0.04)
                    text_frame.margin_top = Inches(0.04)
                    text_frame.margin_bottom = Inches(0.04)
                    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                    paragraph = text_frame.paragraphs[0]
                    paragraph.font.name = item["font_face"]
                    paragraph.font.size = Pt(item["font_size_pt"])
                    paragraph.font.bold = row_index == 0
                    paragraph.font.color.rgb = RGBColor.from_string(
                        item["header_color"] if row_index == 0 else item["body_color"]
                    )
        else:
            shape = slide.shapes.add_textbox(*geometry)
            text_frame = shape.text_frame
            text_frame.clear()
            text_frame.margin_left = 0
            text_frame.margin_right = 0
            text_frame.margin_top = 0
            text_frame.margin_bottom = 0
            text_frame.vertical_anchor = MSO_ANCHOR.TOP
            text_frame.auto_size = (
                MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                if item["fit"] == "shrink"
                else MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
            )
            text_frame.word_wrap = item["word_wrap"]
            paragraph = text_frame.paragraphs[0]
            paragraph.text = item["text"]
            paragraph.alignment = alignments[item["align"]]
            paragraph.space_after = Pt(0)
            paragraph.font.name = item["font_face"]
            paragraph.font.size = Pt(item["font_size_pt"])
            paragraph.font.bold = item["bold"]
            paragraph.font.italic = item["italic"]
            paragraph.font.color.rgb = RGBColor.from_string(item["color"])
        shape.name = item["id"]

    output_path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(output_path)


def _remove_theme_shape_style(shape: object) -> None:
    """Keep scene shapes free from implicit theme effects such as outer shadows."""

    element = shape._element
    for child in list(element):
        if etree.QName(child).localname == "style":
            element.remove(child)


def _add_native_equation(
    slide: object,
    item: dict[str, object],
    geometry: tuple[object, object, object, object],
    base_dir: Path,
    alignments: dict[str, object],
    rgb_color: object,
    vertical_anchor: object,
    points: object,
) -> None:
    """Add editable Office Math with a visual fallback for older viewers."""

    native_shape = slide.shapes.add_textbox(*geometry)
    native_shape.name = str(item["id"])
    text_frame = native_shape.text_frame
    text_frame.clear()
    text_frame.margin_left = 0
    text_frame.margin_right = 0
    text_frame.margin_top = 0
    text_frame.margin_bottom = 0
    text_frame.vertical_anchor = vertical_anchor.MIDDLE
    text_frame.word_wrap = False
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = alignments[str(item["align"])]
    paragraph.space_after = points(0)
    paragraph.font.name = str(item["font_face"])
    paragraph.font.size = points(float(item["font_size_pt"]))
    paragraph.font.color.rgb = rgb_color.from_string(str(item["color"]))
    paragraph_element = paragraph._p
    for child in list(paragraph_element):
        if etree.QName(child).localname != "pPr":
            paragraph_element.remove(child)
    equation_element = etree.fromstring(
        build_text_math_xml(str(item["mathml"]), align=str(item["align"])).encode()
    )
    paragraph_element.append(equation_element)

    fallback_shape = slide.shapes.add_picture(
        _image_source(
            str(item["fallback_src"]),
            base_dir,
            width_in=float(item["w"]),
            height_in=float(item["h"]),
        ),
        *geometry,
    )
    fallback_shape.name = f"{item['id']}.fallback"
    description = str(item["latex"] or "Equation preview")
    fallback_shape._element.nvPicPr.cNvPr.set("descr", description)
    _wrap_equation_fallback(native_shape._element, fallback_shape._element)


def _wrap_equation_fallback(native_element: object, fallback_element: object) -> None:
    """Wrap mutually exclusive native and preview shapes in AlternateContent."""

    parent = native_element.getparent()
    if parent is None or fallback_element.getparent() is not parent:
        raise RuntimeError("equation shapes must share one slide shape tree")
    insertion_index = parent.index(native_element)
    parent.remove(native_element)
    parent.remove(fallback_element)
    alternate = etree.Element(f"{{{MC_NAMESPACE}}}AlternateContent")
    choice = etree.SubElement(
        alternate,
        f"{{{MC_NAMESPACE}}}Choice",
        Requires="a14",
        nsmap={"a14": A14_NAMESPACE},
    )
    choice.append(native_element)
    fallback = etree.SubElement(alternate, f"{{{MC_NAMESPACE}}}Fallback")
    fallback.append(fallback_element)
    parent.insert(insertion_index, alternate)


def _image_source(
    value: str,
    base_dir: Path,
    *,
    width_in: float,
    height_in: float,
) -> Path | BytesIO:
    """Resolve one scene image without network access."""

    mime = ""
    if value.startswith("data:"):
        header, separator, payload = value.partition(",")
        if not separator or ";base64" not in header:
            raise ValueError("PPTX image data URI must be base64 encoded")
        mime = header[5:].split(";", 1)[0].lower()
        raw = base64.b64decode(payload, validate=True)
        return (
            _render_svg(raw, width_in=width_in, height_in=height_in)
            if mime == "image/svg+xml"
            else BytesIO(raw)
        )
    if value.startswith("file://"):
        path = Path(unquote(urlparse(value).path)).resolve(strict=True)
    else:
        path = (base_dir / value).resolve(strict=True)
    if not path.is_file():
        raise ValueError(f"PPTX image is not a regular file: {path}")
    return (
        _render_svg(path.read_bytes(), width_in=width_in, height_in=height_in)
        if path.suffix.lower() == ".svg"
        else path
    )


def _svg_target_pixels(width_in: float, height_in: float) -> tuple[int, int]:
    """Return a 300-PPI placement raster within one bounded memory budget."""

    if (
        not math.isfinite(width_in)
        or not math.isfinite(height_in)
        or width_in <= 0
        or height_in <= 0
    ):
        raise ValueError("SVG placement dimensions must be positive finite inches")
    width = max(1, math.ceil(width_in * _SVG_TARGET_PPI))
    height = max(1, math.ceil(height_in * _SVG_TARGET_PPI))
    limit_scale = min(
        1.0,
        _SVG_MAX_LONG_EDGE_PIXELS / width,
        _SVG_MAX_LONG_EDGE_PIXELS / height,
        math.sqrt(_SVG_MAX_PIXELS / (width * height)),
    )
    if limit_scale < 1:
        width = max(1, math.floor(width * limit_scale))
        height = max(1, math.floor(height * limit_scale))
    return width, height


def _render_svg(raw: bytes, *, width_in: float, height_in: float) -> BytesIO:
    """Rasterize an SVG at its final PowerPoint placement size."""

    try:
        import pymupdf
    except ImportError as exc:
        raise RuntimeError("PyMuPDF is required to export SVG figures") from exc
    document = pymupdf.open(stream=raw, filetype="svg")
    try:
        page = document[0]
        if page.rect.width <= 0 or page.rect.height <= 0:
            raise ValueError("SVG dimensions must be positive")
        target_width, target_height = _svg_target_pixels(width_in, height_in)
        pixmap = page.get_pixmap(
            matrix=pymupdf.Matrix(
                target_width / page.rect.width,
                target_height / page.rect.height,
            ),
            alpha=True,
        )
        return BytesIO(pixmap.tobytes("png"))
    finally:
        document.close()


def _parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("scene", type=Path)
    parser.add_argument("output", type=Path)
    return parser.parse_args()


def main() -> int:
    args = _parse_args()
    try:
        render_scene(args.scene.resolve(strict=True), args.output.resolve())
    except Exception as exc:
        print(str(exc), file=sys.stderr)
        return 1
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
