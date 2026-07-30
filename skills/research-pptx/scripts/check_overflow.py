#!/usr/bin/env python3
"""
Estimate text overflow in PPTX slides BEFORE visual rendering.

Reads every text box, estimates required height based on text content and
font size, and flags boxes where content likely exceeds the container.

This catches overflow that is invisible in rendered images (text is in the
XML but clipped by the box boundary).

Usage:
    python scripts/check_overflow.py output.pptx
    python scripts/check_overflow.py output.pptx --verbose
"""

import argparse
import sys
from pathlib import Path

try:
    from pptx import Presentation
except ImportError:
    print("ERROR: python-pptx required. Install: pip install python-pptx")
    sys.exit(1)


# -- Character width estimation ------------------------------------------------
# Average character width as a fraction of font size (empirical, proportional fonts)
# Constants MUST match generate_slides.js estimateHeight:
#   Latin char width  = font_size_pt * 0.55
#   CJK   char width  = font_size_pt * 1.0
#   Line height       = font_size_pt * 1.4
LATIN_CHAR_WIDTH_FACTOR = 0.55   # e.g., Arial 16pt → avg char ≈ 8.8pt wide
CJK_CHAR_WIDTH_FACTOR = 1.0      # CJK full-width characters
LINE_HEIGHT_FACTOR = 1.4          # line height ≈ 1.4× font size (with leading)
DEFAULT_FONT_SIZE_PT = 16         # fallback when font size is inherited/unset
DEFAULT_MARGIN_INCHES = 0.05      # PptxGenJS default text box internal margin per side


def _is_cjk(ch: str) -> bool:
    """Check if a character is CJK (Chinese/Japanese/Korean)."""
    cp = ord(ch)
    return any([
        0x4E00 <= cp <= 0x9FFF,    # CJK Unified Ideographs
        0x3400 <= cp <= 0x4DBF,    # CJK Extension A
        0xF900 <= cp <= 0xFAFF,    # CJK Compatibility
        0xFF00 <= cp <= 0xFFEF,    # Fullwidth Forms
        0xAC00 <= cp <= 0xD7AF,    # Hangul Syllables
        0x3040 <= cp <= 0x309F,    # Hiragana
        0x30A0 <= cp <= 0x30FF,    # Katakana
    ])


def _estimate_line_count(text: str, font_size_pt: float, box_width_pt: float) -> int:
    """Estimate how many lines a text string needs at given font size and box width.

    Uses the same word-boundary wrapping algorithm as generate_slides.js
    estimateHeight(), so overflow estimates are consistent with the renderer.
    """
    if not text.strip():
        return 1

    latin_w = font_size_pt * LATIN_CHAR_WIDTH_FACTOR
    cjk_w = font_size_pt * CJK_CHAR_WIDTH_FACTOR

    if box_width_pt <= 0:
        return 1

    # Detect CJK-dominant paragraphs (same heuristic as JS renderer)
    cjk_count = sum(1 for ch in text if _is_cjk(ch))
    cjk_dominant = cjk_count > len(text) * 0.3

    if cjk_dominant:
        # Character-based wrapping for CJK
        used = 0.0
        lines = 1
        for ch in text:
            ch_w = cjk_w if _is_cjk(ch) else latin_w
            if used + ch_w > box_width_pt:
                lines += 1
                used = ch_w
            else:
                used += ch_w
    else:
        # Word-boundary wrapping for Latin scripts
        words = text.split()
        used = 0.0
        lines = 1
        space_w = latin_w
        for word in words:
            word_w = sum(cjk_w if _is_cjk(ch) else latin_w for ch in word)
            if used > 0 and used + space_w + word_w > box_width_pt:
                lines += 1
                used = word_w
            else:
                used += (space_w if used > 0 else 0) + word_w

    return lines


def _get_font_size_pt(run) -> float:
    """Get font size in points from a run, with fallback."""
    if run.font.size is not None:
        return run.font.size.pt
    return DEFAULT_FONT_SIZE_PT


def check_text_frame(shape, slide_idx: int) -> list[str]:
    """Check a single text frame for potential overflow. Returns list of warnings."""
    warnings = []

    if not shape.has_text_frame:
        return warnings

    tf = shape.text_frame

    # Box dimensions in points (1 inch = 72pt = 914400 EMU)
    box_w_pt = shape.width / 914400 * 72
    box_h_pt = shape.height / 914400 * 72

    # Account for internal margins
    margin_pt = DEFAULT_MARGIN_INCHES * 72
    usable_w_pt = box_w_pt - 2 * margin_pt
    usable_h_pt = box_h_pt - 2 * margin_pt

    if usable_w_pt <= 0 or usable_h_pt <= 0:
        return warnings

    # Estimate total height needed
    total_height_pt = 0.0
    para_texts = []

    for _para_idx, para in enumerate(tf.paragraphs):
        # Get the effective font size for this paragraph
        font_sizes = [_get_font_size_pt(run) for run in para.runs if run.text.strip()]
        font_size = max(font_sizes) if font_sizes else DEFAULT_FONT_SIZE_PT

        para_text = para.text.strip()
        para_texts.append(para_text)

        # Estimate lines needed
        line_count = _estimate_line_count(para_text, font_size, usable_w_pt)

        # Height for this paragraph
        para_height = line_count * font_size * LINE_HEIGHT_FACTOR
        total_height_pt += para_height

    # Check for overflow
    overflow_pt = total_height_pt - usable_h_pt
    if overflow_pt > 5:  # threshold: 5pt ≈ 0.07 inches
        overflow_in = overflow_pt / 72
        box_x_in = shape.left / 914400
        box_y_in = shape.top / 914400
        box_w_in = shape.width / 914400
        box_h_in = shape.height / 914400

        # Count how many paragraphs are likely clipped
        cumulative = 0.0
        clipped_from = len(para_texts)
        for i, para in enumerate(tf.paragraphs):
            font_sizes = [_get_font_size_pt(r) for r in para.runs if r.text.strip()]
            fs = max(font_sizes) if font_sizes else DEFAULT_FONT_SIZE_PT
            lines = _estimate_line_count(para.text.strip(), fs, usable_w_pt)
            cumulative += lines * fs * LINE_HEIGHT_FACTOR
            if cumulative > usable_h_pt:
                clipped_from = i
                break

        total_paras = len(para_texts)
        clipped_count = total_paras - clipped_from

        severity = "⚠️ WARNING" if overflow_in < 0.5 else "🔴 CRITICAL"

        warnings.append(
            f"{severity}: Slide {slide_idx}, text box at "
            f"({box_x_in:.1f}\", {box_y_in:.1f}\") size {box_w_in:.1f}\"×{box_h_in:.1f}\":\n"
            f"    Estimated text height: {total_height_pt / 72:.2f}\" "
            f"(box usable: {usable_h_pt / 72:.2f}\") — "
            f"overflow by {overflow_in:.2f}\"\n"
            f"    {total_paras} paragraphs total, ~{clipped_count} likely clipped\n"
            f"    Last visible paragraph: \"{para_texts[clipped_from][:60]}...\""
            if clipped_from < len(para_texts) and len(para_texts[clipped_from]) > 60
            else f"    {total_paras} paragraphs total, ~{clipped_count} likely clipped"
        )

    return warnings


def check_table_overflow(shape, slide_idx: int) -> list[str]:
    """Check if table cells or overall table dimensions overflow."""
    warnings = []
    if not shape.has_table:
        return warnings

    table = shape.table
    total_rows = len(table.rows)

    # Shape dimensions
    shape_h_pt = shape.height / 914400 * 72
    shape_w_pt = shape.width / 914400 * 72
    usable_shape_h = shape_h_pt - 2 * DEFAULT_MARGIN_INCHES * 72

    cell_margin_pt = DEFAULT_MARGIN_INCHES * 72
    fallback_col_w = shape_w_pt / max(len(table.columns), 1)
    column_widths_pt = [
        (column.width / 914400 * 72) if column.width else fallback_col_w
        for column in table.columns
    ]

    # Estimate per-row text height
    total_text_h_pt = 0.0
    for row_idx, row in enumerate(table.rows):
        row_text_h = 0.0
        # Estimate row height: 0.45" header + 0.4" per data row
        nominal_h_pt = (0.45 * 72) if row_idx == 0 else (0.4 * 72)
        usable_row_h = nominal_h_pt - 2 * cell_margin_pt

        for col_idx, cell in enumerate(row.cells):
            cell_text = cell.text.strip()
            if not cell_text:
                continue
            cell_w_pt = (
                column_widths_pt[col_idx]
                if col_idx < len(column_widths_pt)
                else fallback_col_w
            )
            usable_cell_w = max(cell_w_pt - 2 * cell_margin_pt, 1.0)
            font_size = DEFAULT_FONT_SIZE_PT
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size:
                        font_size = run.font.size.pt
                        break

            para_h = 0.0
            for para in cell.text_frame.paragraphs:
                para_text = para.text.strip()
                if not para_text:
                    para_h += font_size * LINE_HEIGHT_FACTOR
                    continue
                lines = _estimate_line_count(para_text, font_size, usable_cell_w)
                para_h += lines * font_size * LINE_HEIGHT_FACTOR

            row_text_h = max(row_text_h, para_h)

        total_text_h_pt += max(row_text_h, usable_row_h)

    overflow_pt = total_text_h_pt - usable_shape_h
    if overflow_pt > 10:
        overflow_in = overflow_pt / 72
        severity = "🔴 CRITICAL" if overflow_in >= 0.5 else "⚠️ WARNING"
        warnings.append(
            f"{severity}: Slide {slide_idx}, table: estimated content "
            f"{total_text_h_pt / 72:.2f}\" exceeds shape {usable_shape_h / 72:.2f}\" "
            f"by {overflow_in:.2f}\" ({total_rows} rows)"
        )

    # Also flag individual cells with very long wrapping text
    for row_idx, row in enumerate(table.rows):
        for col_idx, cell in enumerate(row.cells):
            cell_text = cell.text.strip()
            if not cell_text:
                continue
            cell_w_pt = (
                column_widths_pt[col_idx]
                if col_idx < len(column_widths_pt)
                else fallback_col_w
            )
            usable_cell_w = max(cell_w_pt - 2 * cell_margin_pt, 1.0)
            font_size = DEFAULT_FONT_SIZE_PT
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size:
                        font_size = run.font.size.pt
                        break
            lines_needed = _estimate_line_count(cell_text, font_size, usable_cell_w)
            if lines_needed > 3:
                warnings.append(
                    f"⚠️ WARNING: Slide {slide_idx}, table row {row_idx + 1} "
                    f"col {col_idx + 1}: cell text may wrap to {lines_needed} lines "
                    f"(\"{cell_text[:40]}...\")"
                )

    return warnings


def main():
    parser = argparse.ArgumentParser(description="Check PPTX for potential text overflow")
    parser.add_argument("pptx_path", help="Path to .pptx file")
    parser.add_argument("--verbose", "-v", action="store_true", help="Show per-slide summary")
    args = parser.parse_args()

    pptx_path = Path(args.pptx_path)
    if not pptx_path.exists():
        print(f"File not found: {pptx_path}")
        sys.exit(1)

    prs = Presentation(str(pptx_path))
    all_warnings = []

    for slide_idx, slide in enumerate(prs.slides, 1):
        slide_warnings = []
        for shape in slide.shapes:
            slide_warnings.extend(check_text_frame(shape, slide_idx))
            slide_warnings.extend(check_table_overflow(shape, slide_idx))

        if args.verbose and not slide_warnings:
            print(f"  Slide {slide_idx}: OK")

        all_warnings.extend(slide_warnings)

    if all_warnings:
        print(f"\nFound {len(all_warnings)} potential overflow(s):\n")
        for w in all_warnings:
            print(w)
            print()
        sys.exit(1)
    else:
        print("✅ No text overflow detected.")
        sys.exit(0)


if __name__ == "__main__":
    main()
