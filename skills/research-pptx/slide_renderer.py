"""PPTX generation via Node.js + PptxGenJS.

Writes slide data as JSON, invokes generate_slides.js asynchronously,
returns PPTX path.
"""

from __future__ import annotations

import asyncio
import importlib.util as _ilu
import json
import logging
import os
import shutil
import sys as _sys
from pathlib import Path
from typing import Any

_spec = _ilu.spec_from_file_location(
    "research_pptx_models", Path(__file__).resolve().parent / "models.py"
)
_models = _sys.modules.get("research_pptx_models")
if _models is None:
    _models = _ilu.module_from_spec(_spec)
    _sys.modules["research_pptx_models"] = _models
    _spec.loader.exec_module(_models)

PresentationPlan = _models.PresentationPlan

logger = logging.getLogger(__name__)

_SCRIPTS_DIR = Path(__file__).resolve().parent / "scripts"


class RendererDependencyError(RuntimeError):
    """A setup-time renderer dependency is unavailable."""

    def __init__(self, message: str, *, code: str, missing: list[str]) -> None:
        super().__init__(message)
        self.code = code
        self.missing = missing


def preflight_renderer(node_runtime_dir: str | Path | None = None) -> str:
    """Validate the renderer without downloading or mutating dependencies."""
    node = shutil.which("node") or shutil.which("node.exe")
    if not node:
        raise RendererDependencyError(
            "research-pptx renderer is not ready; install Node.js >= 20.9",
            code="node_unavailable",
            missing=["node"],
        )

    runtime_dir = Path(node_runtime_dir) if node_runtime_dir is not None else _SCRIPTS_DIR
    node_modules = runtime_dir / "node_modules"
    missing_packages = [
        package
        for package in ("pptxgenjs", "sharp")
        if not (node_modules / package).is_dir()
    ]
    if missing_packages:
        raise RendererDependencyError(
            "research-pptx Node packages are not pre-installed: "
            + ", ".join(missing_packages)
            + ". Complete renderer setup from the host CLI before retrying.",
            code="runtime_dependency_missing",
            missing=missing_packages,
        )
    return node


def resolve_figure_paths(
    plan: PresentationPlan,
    figures: list[dict[str, str]],
    work_dir: str,
    *,
    placeholder_mode: str = "auto",
) -> PresentationPlan:
    """Replace figure_N references with actual file paths.

    ``placeholder_mode``:
      - ``"auto"`` (default): if the referenced figure cannot be resolved AND
        we can generate a placeholder image, keep the slide type and stamp a
        "Figure N — placeholder" image so downstream review can spot missing
        assets and the layout stays visually correct.
      - ``"downgrade"``: legacy behaviour — silently downgrade
        ``content_figure`` → ``content`` and drop the figure caption.

    A resolvable reference always wins over a placeholder; the placeholder is
    only synthesised when the reference literally has no on-disk file. This
    matches the figure placeholder requirement: users regularly ask for a deck first
    and swap in the real figures later.
    """
    n = len(figures)
    downgraded = 0
    placeheld = 0

    for slide in plan.slides:
        fp = slide.figure_path
        if not fp:
            continue

        resolved: str | None = None
        placeholder_hint = ""

        if fp == "__placeholder__":
            hint = slide.figure_caption or "figure to be provided"
            ph_path = _render_figure_placeholder(work_dir, hint)
            if ph_path is not None:
                slide.figure_path = ph_path
                if slide.figure_caption:
                    if "[placeholder]" not in slide.figure_caption.lower():
                        slide.figure_caption = (
                            f"[placeholder needed] {slide.figure_caption}"
                        )
                else:
                    slide.figure_caption = f"[placeholder needed] {hint[:80]}"
                placeheld += 1
                continue
            # Pillow unavailable — fall through to the downgrade branch.
            slide.figure_path = None
            slide.figure_caption = None
            if slide.slide_type == "content_figure":
                slide.slide_type = "content"
                downgraded += 1
            elif slide.slide_type == "full_figure":
                slide.slide_type = "content"
                if not slide.bullets:
                    slide.bullets = [slide.title or "See notes for details"]
                downgraded += 1
            continue

            # Pattern: figure_N
        if fp.startswith("figure_"):
            try:
                idx = int(fp.split("_", 1)[1])
                if 0 <= idx < n:
                    cand = figures[idx].get("path", "")
                    if cand and os.path.exists(cand):
                        resolved = cand
                    # If the plan references figure_N but no such figure
                    # exists on disk, remember the caption so the placeholder
                    # can be labelled meaningfully.
                    if not resolved and 0 <= idx < n:
                        placeholder_hint = (
                            figures[idx].get("caption", "") or f"Figure {idx}"
                        )
                    else:
                        placeholder_hint = f"Figure {idx}"
            except (ValueError, IndexError):
                pass
        # Pattern: absolute or relative path
        elif os.path.exists(fp):
            resolved = fp
        elif os.path.exists(os.path.join(work_dir, fp)):
            resolved = os.path.join(work_dir, fp)

        if resolved:
            slide.figure_path = resolved
            continue

        # ── Placeholder path ──
        # Try to render a lightweight PNG placeholder so the layout stays
        # intact. If placeholder generation fails we fall back to the legacy
        # downgrade behaviour.
        if placeholder_mode == "auto":
            ph_label = placeholder_hint or slide.figure_caption or fp
            ph_path = _render_figure_placeholder(work_dir, ph_label)
            if ph_path is not None:
                slide.figure_path = ph_path
                # Keep the caption but mark it visually so reviewers see it.
                if slide.figure_caption:
                    if "[placeholder]" not in slide.figure_caption.lower():
                        slide.figure_caption = (
                            f"[placeholder] {slide.figure_caption}"
                        )
                else:
                    slide.figure_caption = f"[placeholder] {ph_label[:80]}"
                placeheld += 1
                continue

        logger.warning(
            "[slide-renderer] Unresolvable figure reference '%s' on slide "
            "'%s' (type=%s) — downgrading layout",
            fp, (slide.title or "")[:40], slide.slide_type,
        )
        slide.figure_path = None
        slide.figure_caption = None
        if slide.slide_type == "content_figure":
            slide.slide_type = "content"
            downgraded += 1
        elif slide.slide_type == "full_figure":
            slide.slide_type = "content"
            if not slide.bullets:
                slide.bullets = [slide.title or "See paper for details"]
            downgraded += 1

    if placeheld:
        logger.info(
            "[slide-renderer] Generated %d figure placeholder(s)", placeheld,
        )
    if downgraded:
        logger.info(
            "[slide-renderer] Downgraded %d slides due to missing figures",
            downgraded,
        )
    return plan


def _render_figure_placeholder(work_dir: str, label: str) -> str | None:
    """Render a grey "figure placeholder" PNG so a slide keeps its layout
    even when the real asset is missing.

    Uses Pillow (already a transitive dep via matplotlib / pymupdf). Returns
    the placeholder file path, or None if Pillow is unavailable.
    """
    try:
        from PIL import Image, ImageDraw, ImageFont
    except ImportError:
        return None

    ph_dir = Path(work_dir) / "placeholders"
    ph_dir.mkdir(parents=True, exist_ok=True)
    # Stable name from label so a re-render reuses the same file.
    import hashlib as _h
    fname = "ph_" + _h.sha1(label.encode("utf-8")).hexdigest()[:10] + ".png"
    out = ph_dir / fname
    if out.is_file():
        return str(out)

    # 4:3 canvas, projection-friendly grey.
    W, H = 1200, 900
    img = Image.new("RGB", (W, H), color=(235, 238, 244))
    draw = ImageDraw.Draw(img)

    # Border + diagonal cross (the classic "missing image" look).
    draw.rectangle([(6, 6), (W - 6, H - 6)], outline=(180, 190, 210), width=4)
    draw.line([(6, 6), (W - 6, H - 6)], fill=(200, 210, 225), width=3)
    draw.line([(W - 6, 6), (6, H - 6)], fill=(200, 210, 225), width=3)

    text_lines = ["Figure placeholder", ""]
    remaining = (label or "").strip()
    for _ in range(4):
        if not remaining:
            break
        text_lines.append(remaining[:40])
        remaining = remaining[40:]

    try:
        font = ImageFont.truetype("arial.ttf", 42)
        small = ImageFont.truetype("arial.ttf", 28)
    except OSError:
        font = ImageFont.load_default()
        small = font

    # Centre the text block.
    y = H // 2 - (len(text_lines) * 40) // 2
    for i, line in enumerate(text_lines):
        f = font if i == 0 else small
        try:
            bbox = draw.textbbox((0, 0), line, font=f)
            tw = bbox[2] - bbox[0]
        except AttributeError:
            tw = draw.textsize(line, font=f)[0]
        draw.text(((W - tw) // 2, y), line, fill=(90, 100, 120), font=f)
        y += 55 if i == 0 else 40

    img.save(out, "PNG")
    return str(out)

async def render_pptx(
    plan: PresentationPlan,
    work_dir: str,
    output_filename: str = "presentation.pptx",
    *,
    node_runtime_dir: str | Path | None = None,
) -> str:
    """Generate PPTX by calling Node.js script asynchronously.

    Returns the path to the generated .pptx file.
    """
    node = preflight_renderer(node_runtime_dir)
    js_script = _SCRIPTS_DIR / "generate_slides.js"

    if not js_script.exists():
        raise FileNotFoundError(f"JS renderer not found: {js_script}")

    for s in plan.slides:
        if s.figure_caption:
            cap = s.figure_caption.strip()
            if len(cap) > 100:
                cap = cap[:97] + "..."
            s.figure_caption = cap

    # Write slide data JSON
    slide_data = {
        "config": {
            "title": plan.title,
            "authors": plan.authors,
            "affiliation": plan.affiliation,
            "venue": plan.venue,
            "colors": plan.color_theme,
            "headerFont": getattr(plan, "header_font", "") or "Arial Black",
            "bodyFont": getattr(plan, "body_font", "") or "Arial",
            # structured template master (background/logo/title-band).
            "templateMaster": getattr(plan, "template_master", {}) or {},
        },
        "slides": [s.model_dump() for s in plan.slides],
    }

    data_path = os.path.join(work_dir, "slide_data.json")
    with open(data_path, "w", encoding="utf-8") as f:
        json.dump(slide_data, f, ensure_ascii=False, indent=2)

    output_path = os.path.join(work_dir, output_filename)

    # Run Node.js asynchronously
    runtime_dir = Path(node_runtime_dir) if node_runtime_dir is not None else _SCRIPTS_DIR
    node_path = str(runtime_dir / "node_modules")
    existing_node_path = os.environ.get("NODE_PATH", "").strip()
    if existing_node_path:
        node_path = os.pathsep.join((node_path, existing_node_path))
    node_env = {**os.environ, "NODE_PATH": node_path}
    proc = await asyncio.create_subprocess_exec(
        node, str(js_script), data_path, output_path,
        stdout=asyncio.subprocess.PIPE,
        stderr=asyncio.subprocess.PIPE,
        cwd=work_dir,
        env=node_env,
    )

    try:
        stdout, stderr = await asyncio.wait_for(proc.communicate(), timeout=120)
    except TimeoutError as exc:
        proc.kill()
        raise RuntimeError("PptxGenJS generation timed out after 120s") from exc

    if proc.returncode != 0:
        err_text = stderr.decode(errors="replace")[:2000]
        logger.error("[slide-renderer] Node.js stderr: %s", err_text)
        raise RuntimeError(f"PptxGenJS generation failed: {err_text[:500]}")

    if not os.path.exists(output_path):
        raise RuntimeError(f"PPTX file not created at {output_path}")

    logger.info("[slide-renderer] Generated: %s", output_path)
    return output_path


async def run_overflow_check(pptx_path: str) -> list[str]:
    """Run overflow detection asynchronously. Returns list of warning strings."""
    check_script = _SCRIPTS_DIR / "check_overflow.py"
    if not check_script.exists():
        logger.warning("[slide-renderer] check_overflow.py not found, skipping QA")
        return []

    try:
        proc = await asyncio.create_subprocess_exec(
            "python", str(check_script), pptx_path,
            stdout=asyncio.subprocess.PIPE,
            stderr=asyncio.subprocess.PIPE,
        )
        stdout, stderr = await asyncio.wait_for(proc.communicate(), timeout=30)

        if proc.returncode != 0:
            return [
                line for line in stdout.decode(errors="replace").split("\n")
                if line.strip()
            ]
        return []
    except TimeoutError:
        logger.warning("[slide-renderer] Overflow check timed out")
        return []
    except Exception as e:
        logger.warning("[slide-renderer] Overflow check failed: %s", e)
        return []

def _is_cjk_char(ch: str) -> bool:
    cp = ord(ch)
    return any([
        0x4E00 <= cp <= 0x9FFF,
        0x3400 <= cp <= 0x4DBF,
        0xF900 <= cp <= 0xFAFF,
        0xFF00 <= cp <= 0xFFEF,
        0xAC00 <= cp <= 0xD7AF,
        0x3040 <= cp <= 0x309F,
        0x30A0 <= cp <= 0x30FF,
    ])


def _estimate_text_height_pt(
    text: str, font_size_pt: float, box_width_pt: float,
) -> float:
    """Estimate rendered text height in points.

    Parameters aligned with generate_slides.js ``estimateHeight`` so the
    overflow check uses the same constants as the renderer:
      Latin char width  = font_size_pt * 0.55  (was 0.52)
      CJK   char width  = font_size_pt * 1.0
      Line height       = font_size_pt * 1.4   (was 1.35)
    """
    if not text.strip():
        return font_size_pt * 1.4

    latin_w = font_size_pt * 0.55
    cjk_w = font_size_pt * 1.0
    line_h = font_size_pt * 1.4

    # Detect CJK-dominant paragraphs (same heuristic as JS renderer)
    cjk_count = sum(1 for ch in text if _is_cjk_char(ch))
    cjk_dominant = cjk_count > len(text) * 0.3

    if cjk_dominant:
        # Character-based wrapping
        used = 0.0
        lines = 1
        for ch in text:
            ch_w = cjk_w if _is_cjk_char(ch) else latin_w
            if used + ch_w > box_width_pt:
                lines += 1
                used = ch_w
            else:
                used += ch_w
    else:
        # Word-boundary wrapping for Latin scripts
        words = text.split()
        used = 0.0
        lines = 1
        space_w = latin_w
        for word in words:
            word_w = sum(cjk_w if _is_cjk_char(ch) else latin_w for ch in word)
            if used > 0 and used + space_w + word_w > box_width_pt:
                lines += 1
                used = word_w
            else:
                used += (space_w if used > 0 else 0) + word_w

    return lines * line_h + 0.25  # +0.25 safety pad matches JS


def _check_table_overflow_structured(
    shape, default_font_pt: float, margin_in: float,
) -> float:
    """Estimate max text overflow (in inches) across all cells of a table shape.

    Uses the same _estimate_text_height_pt as text-frame checking so constants
    stay consistent with the JS renderer.
    """
    if not shape.has_table:
        return 0.0

    table = shape.table
    total_rows = len(table.rows)
    total_cols = len(table.columns)

    # Shape dimensions in points
    shape_h_pt = shape.height / 914400 * 72
    shape_w_pt = shape.width / 914400 * 72
    margin_pt = margin_in * 72
    usable_shape_h = shape_h_pt - 2 * margin_pt
    usable_shape_w = shape_w_pt - 2 * margin_pt

    if usable_shape_h <= 0 or usable_shape_w <= 0:
        return 0.0

    # Estimate row height from the shape: even distribution as first guess,
    # then read actual row heights from the XML if available.
    row_heights_pt: list[float] = []
    for row in table.rows:
        rh = row.height
        if rh is not None and rh > 0:
            row_heights_pt.append(rh / 914400 * 72)
    if not row_heights_pt:
        # Fallback: JS renderer uses 0.45" header + 0.4"/data row
        if total_rows > 0:
            row_heights_pt = [0.45 * 72] + [0.4 * 72] * (total_rows - 1)

    if not row_heights_pt:
        return 0.0

    total_text_h_pt = 0.0
    cell_margin_pt = margin_pt  # cells have their own internal margins

    fallback_col_w = shape_w_pt / max(total_cols, 1)
    column_widths_pt = [
        (column.width / 914400 * 72) if column.width else fallback_col_w
        for column in table.columns
    ]

    for row_idx, row in enumerate(table.rows):
        row_text_h = 0.0
        row_height_pt = (
            row_heights_pt[row_idx]
            if row_idx < len(row_heights_pt)
            else row_heights_pt[-1]
        )
        usable_row_h = row_height_pt - 2 * cell_margin_pt

        for col_idx, cell in enumerate(row.cells):
            cell_w_pt = (
                column_widths_pt[col_idx]
                if col_idx < len(column_widths_pt)
                else fallback_col_w
            )
            usable_cell_w = max(cell_w_pt - 2 * cell_margin_pt, 1.0)
            cell_para_h = 0.0
            for para in cell.text_frame.paragraphs:
                font_sizes = []
                for run in para.runs:
                    if run.text.strip() and run.font.size is not None:
                        font_sizes.append(run.font.size.pt)
                fs = max(font_sizes) if font_sizes else default_font_pt
                cell_para_h += _estimate_text_height_pt(
                    para.text.strip(), fs, usable_cell_w,
                )
            row_text_h = max(row_text_h, cell_para_h)

        # Accumulate: if a row's text exceeds its nominal height,
        # the excess pushes *this* row down and steals space from
        # rows below it, so we add the delta to total_text_h.
        allocated = max(usable_row_h, 0)
        total_text_h_pt += max(row_text_h, allocated)

    overflow_pt = total_text_h_pt - usable_shape_h
    if overflow_pt > 5:
        return overflow_pt / 72
    return 0.0


def check_overflow_structured(pptx_path: str) -> list[dict[str, Any]]:
    """In-process overflow check returning structured per-slide results.

    Returns a list of dicts:
      [{"slide_index": 0, "overflow_inches": 0.42, "severity": "warning"}, ...]

    Severity: "warning" (<0.5"), "critical" (≥0.5").
    Empty list means no overflow detected.
    """
    try:
        from pptx import Presentation as PptxPresentation
    except ImportError:
        logger.warning("python-pptx not available; skipping structured overflow check")
        return []

    _DEFAULT_FONT_PT = 16.0
    _MARGIN_IN = 0.05

    try:
        prs = PptxPresentation(str(pptx_path))
    except Exception:
        logger.debug("Failed to open PPTX for overflow check", exc_info=True)
        return []

    results: list[dict[str, Any]] = []

    for slide_idx, slide in enumerate(prs.slides):
        max_overflow_in = 0.0

        for shape in slide.shapes:
            # ── Text frame overflow check ──
            if shape.has_text_frame:
                box_w_pt = shape.width / 914400 * 72
                box_h_pt = shape.height / 914400 * 72
                margin_pt = _MARGIN_IN * 72
                usable_w = box_w_pt - 2 * margin_pt
                usable_h = box_h_pt - 2 * margin_pt

                if usable_w <= 0 or usable_h <= 0:
                    continue

                total_h = 0.0
                for para in shape.text_frame.paragraphs:
                    font_sizes = []
                    for run in para.runs:
                        if run.text.strip() and run.font.size is not None:
                            font_sizes.append(run.font.size.pt)
                    fs = max(font_sizes) if font_sizes else _DEFAULT_FONT_PT
                    total_h += _estimate_text_height_pt(para.text.strip(), fs, usable_w)

                overflow_pt = total_h - usable_h
                if overflow_pt > 5:
                    max_overflow_in = max(max_overflow_in, overflow_pt / 72)

            # ── Table overflow check ──
            if shape.has_table:
                max_overflow_in = max(
                    max_overflow_in,
                    _check_table_overflow_structured(shape, _DEFAULT_FONT_PT, _MARGIN_IN),
                )

        if max_overflow_in > 0:
            results.append({
                "slide_index": slide_idx,
                "overflow_inches": round(max_overflow_in, 2),
                "severity": "critical" if max_overflow_in >= 0.5 else "warning",
            })

    return results

def build_layout_report(
    plan: PresentationPlan,
    overflow: list[dict[str, Any]],
) -> list[dict[str, Any]]:
    """Translate the rendered deck into a *textual* visual report.

    DeepSeek (text-only) cannot see the rendered PNG, so we surface the layout
    metrics we already compute (figure aspect/layout mode, bullet lengths,
    fill ratio, figure↔section match) as structured text. This is the
    'structured visual feedback' layer the LLM critique reads instead of pixels.
    """
    overflow_by_idx = {o["slide_index"]: o for o in overflow}
    report: list[dict[str, Any]] = []
    for i, s in enumerate(plan.slides):
        row: dict[str, Any] = {
            "slide_index": i,
            "slide_type": s.slide_type,
            "title": (s.title or "")[:100],
            "warnings": [],
        }
        if s.bullets:
            lens = [len(b) for b in s.bullets]
            row["bullet_count"] = len(s.bullets)
            row["max_bullet_chars"] = max(lens)
            row["avg_bullet_chars"] = round(sum(lens) / len(lens), 1)
            if max(lens) > 70:
                row["warnings"].append("bullet_too_long")
            if len(s.bullets) < 2 and s.slide_type in ("content", "content_figure"):
                row["warnings"].append("too_sparse")
        if s.figure_path:
            has_fig = os.path.exists(s.figure_path)
            row["figure"] = {
                "assigned": os.path.basename(s.figure_path) if has_fig else None,
                "resolved": has_fig,
                "caption": (s.figure_caption or "")[:120],
            }
            if not has_fig:
                row["warnings"].append("figure_missing")
        if s.slide_type == "table":
            row["table_rows"] = len(s.table_rows or [])
            row["table_cols"] = len(s.table_headers or [])
            if not s.table_rows or not s.table_headers:
                row["warnings"].append("empty_table")
        if i in overflow_by_idx:
            o = overflow_by_idx[i]
            row["overflow_inches"] = o["overflow_inches"]
            row["warnings"].append(f"overflow_{o['severity']}")
        report.append(row)
    return report
