# template_backend.py
"""Template-faithful PPTX rendering by REUSING a user template as container.

Opens the user's PPTX template with python-pptx, adds slides using the
template's own slideLayouts, and fills content into placeholders. For slide
types that need custom layout beyond what placeholder offers (two_column,
icon_rows, steps, emphasis_box), we add absolutely-positioned shapes on top
of a clean base layout.
"""
from __future__ import annotations

import importlib.util as _ilu
import logging
import math
import os
import sys as _sys
from copy import deepcopy
from pathlib import Path
from typing import Any

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

logger = logging.getLogger(__name__)

# ── Placeholder type → semantic role mapping ──
_PH_ROLE_MAP = {
    "title": "title",
    "ctrTitle": "title",
    "subTitle": "subtitle",
    "body": "body",
    "pic": "picture",
    "dt": "date",
    "ftr": "footer",
    "sldNum": "slideNumber",
}

# Default color fallbacks (template theme colors are preferred)
_DEFAULT_COLORS = {
    "primary": "1E2761",
    "secondary": "CADCFC",
    "accent": "FFFFFF",
    "dark": "0F1535",
    "bodyText": "2D2D2D",
    "muted": "6B7280",
    "tableFill": "F0F4FF",
    "tableHead": "1E2761",
}

# ── Layout-name → role heuristics (Chinese + English, first match wins) ──
_LAYOUT_NAME_HINTS: tuple[tuple[str, tuple[str, ...]], ...] = (
    ("toc",         ("\u76ee\u5f55", "\u5927\u7eb2", "\u7ae0\u8282\u5217\u8868", "toc",
                     "contents", "agenda", "outline")),
    ("closing",     ("\u7ed3\u675f", "\u8c22\u8c22", "\u611f\u8c22", "\u81f4\u8c22", "\u95ee\u7b54",
                     "closing", "thank", "q&a", "questions", "end")),
    ("section",     ("\u7ae0\u8282", "\u8fc7\u6e21", "\u5206\u5272", "\u5206\u9875", "\u7ae0",
                     "section", "chapter", "divider", "transition")),
    ("title",       ("\u6807\u9898\u9875", "\u5c01\u9762", "\u9996\u9875",
                     "title slide", "cover", "opening")),
    ("two_content", ("\u4e24\u680f", "\u53cc\u680f", "\u5bf9\u6bd4", "\u5e76\u5217",
                     "two content", "two column", "comparison", "compare")),
    ("picture",     ("\u56fe\u6587", "\u56fe\u7247", "\u914d\u56fe",
                     "picture", "image with", "figure")),
    ("content",     ("\u5185\u5bb9", "\u6b63\u6587",
                     "content", "body", "text")),
    ("blank",       ("\u7a7a\u767d", "blank", "empty")),
)

_SKILL_DIR = Path(__file__).resolve().parent

# ── Dynamic slide dimensions (templates may be 16:9 wide 13.333x7.5, 4:3, custom) ──
def _slide_dims(slide) -> tuple[float, float]:
    """Return (width_in, height_in) of the owning presentation."""
    try:
        prs = slide.part.package.presentation_part.presentation
        return prs.slide_width / 914400, prs.slide_height / 914400
    except AttributeError:
        return 10.0, 5.625  # standard 16:9 fallback

def _content_region(slide) -> tuple[float, float, float, float]:
    """Return (x, y, w, h) of the usable body content region in inches.
    Derives from the actual slide dimensions so both 10x5.625 (4:3-ish) and
    13.333x7.5 (16:9 wide) templates place content correctly."""
    W, H = _slide_dims(slide)
    return _MARGIN_L, _BODY_TOP, W - _MARGIN_L - _MARGIN_R, H - _BODY_TOP - _BOTTOM

# Content-region margins used by all absolute-positioning fallbacks.
# Chosen to align with typical template placeholder geometry.
_MARGIN_L = 0.9
_MARGIN_R = 0.9
_BODY_TOP = 2.0    # aligns with template body/object placeholder top
_BOTTOM   = 0.6

# Content placeholder types (BODY = classic body, OBJECT = "content" placeholder).
_CONTENT_PH_TYPES = {"BODY", "OBJECT", "body", "object"}

def _load_sibling(mod_name: str):
    """Load a sibling module by file path (same as engine.py's helper)."""
    cached = _sys.modules.get(f"research_pptx_{mod_name}")
    if cached is not None:
        return cached
    spec = _ilu.spec_from_file_location(
        f"research_pptx_{mod_name}", _SKILL_DIR / f"{mod_name}.py"
    )
    module = _ilu.module_from_spec(spec)
    _sys.modules[f"research_pptx_{mod_name}"] = module
    spec.loader.exec_module(module)
    return module

def _new_textbox(slide, x, y, w, h):
    """Create a textbox with word wrap enabled by default.

    Every caller passes ``Inches(...)`` (i.e. already-wrapped Emu integers),
    so this helper accepts pre-wrapped values and forwards them straight to
    ``add_textbox``. The previous version re-wrapped each arg via
    ``Inches(...)``, which double-converted valid Emu values into
    astronomical dimensions and — because callers also omitted the slide —
    raised ``missing 1 required positional argument: 'h'``.

    Forcing ``word_wrap=True`` here is the whole reason the wrapper exists:
    python-pptx leaves ``word_wrap`` as None for user-created textboxes,
    which most viewers treat as "don't wrap" and clip long CJK titles /
    captions / bullets on the right edge.
    """
    tb = slide.shapes.add_textbox(x, y, w, h)
    tb.text_frame.word_wrap = True
    return tb

def _hex_to_rgb(hex_str: str) -> RGBColor:
    h = hex_str.lstrip("#")
    if len(h) == 6:
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
    return RGBColor(0x1E, 0x27, 0x61)


def _role_from_layout_name(name: str) -> str | None:
    """Return a role for *name* by keyword match, or None if nothing matches."""
    if not name:
        return None
    lower = name.strip().lower()
    for role, keywords in _LAYOUT_NAME_HINTS:
        for kw in keywords:
            if kw in lower or kw in name:  # exact-case match for CJK
                return role
    return None

async def classify_layouts_with_llm(
    llm_client, layouts_info: list[dict],
) -> dict[str, int]:
    """Ask the LLM to map semantic slide roles to layout indices.

    This is the "let the LLM decide how to use each template page" hook.
    Layout NAMES are the primary signal (they carry the designer's intent);
    placeholder structure is passed as a secondary hint so the model can
    disambiguate lookalikes.

    Returns a partial {role: index} mapping. The caller MERGES this on top
    of the heuristic result, so LLM decisions override the heuristics but
    the heuristics still fill in anything the LLM omits.

    Returns {} on any failure — the render path stays alive on heuristics.
    """
    import json
    import re

    if not layouts_info or llm_client is None:
        return {}

    lines: list[str] = []
    for layout_info in layouts_info:
        ph_types = [p.get("type", "") for p in layout_info.get("placeholders", [])]
        lines.append(
            f"  {layout_info['index']}: name={layout_info.get('name', '')!r}, "
            f"placeholders={ph_types}"
        )

    system = (
        "You classify PPTX slide layouts by semantic role. "
        "Output only JSON. Never invent layout indices."
    )
    user = (
        "Here is a template's layout inventory. Assign each role below to "
        "the BEST matching layout index (or omit that role if no layout "
        "fits well). Layout names may be in Chinese or English.\n\n"
        f"Layouts:\n{chr(10).join(lines)}\n\n"
        "Roles:\n"
        "  title        — cover / opening slide\n"
        "  toc          — table of contents / agenda\n"
        "  section      — chapter / section divider (usually big title only)\n"
        "  content      — standard title + body page\n"
        "  two_content  — two side-by-side content areas\n"
        "  picture      — title + one large picture area\n"
        "  closing      — thank-you / Q&A / ending page\n"
        "  blank        — fully empty layout\n\n"
        "Rules:\n"
        "  - Prefer the layout NAME as evidence; use placeholders to "
        "disambiguate lookalikes.\n"
        "  - The same index may back multiple roles (e.g., blank fallback).\n"
        "  - Omit a role entirely if no layout fits.\n\n"
        'Return JSON only, e.g. {"title": 0, "toc": 1, "content": 3, '
        '"closing": 6}'
    )

    try:
        raw = await llm_client.chat(
            system, user, temperature=0.1, max_tokens=512,
        )
    except Exception:
        return {}

    if not raw:
        return {}

    # Accept fenced, unfenced, or embedded JSON.
    candidates = [raw, re.sub(r"^```[a-zA-Z]*\n|\n```$", "", raw.strip())]
    m = re.search(r"\{[^{}]*\}", raw, re.DOTALL)
    if m:
        candidates.append(m.group(0))

    for cand in candidates:
        try:
            obj = json.loads(cand)
        except (TypeError, json.JSONDecodeError):
            continue
        if not isinstance(obj, dict):
            continue
        n = len(layouts_info)
        out: dict[str, int] = {}
        for k, v in obj.items():
            if not isinstance(k, str):
                continue
            try:
                idx = int(v)
            except (TypeError, ValueError):
                continue
            if 0 <= idx < n:
                out[k] = idx
        return out
    return {}

def can_reuse_template(template_path: str) -> dict[str, Any] | None:
    """Inspect template PPTX and return layout info if reuse is viable.

    Role classification runs in two passes:
      1. Layout name keywords (Chinese + English).  ← most reliable
      2. Placeholder-shape heuristics (legacy).      ← fallback

    Both write into the same `layout_roles` dict; the first pass sets a
    role only if it's still empty, so name-based hits win.
    """
    if not template_path or not os.path.exists(template_path):
        return None
    try:
        prs = Presentation(template_path)
    except Exception:
        return None
    if not prs.slide_layouts:
        return None

    layouts_info: list[dict[str, Any]] = []
    layout_roles: dict[str, int] = {}

    # ── Pass 0: collect layout info once. ──
    for i, layout in enumerate(prs.slide_layouts):
        ph_list: list[dict[str, Any]] = []
        for ph in layout.placeholders:
            if ph.placeholder_format.type is not None:
                # e.g. "CENTER_TITLE (15)" → "CENTER_TITLE"
                ph_type = (
                    str(ph.placeholder_format.type)
                    .split(".")[-1].split(" ")[0]
                )
            else:
                ph_type = "unknown"
            ph_list.append({
                "type": ph_type,
                "idx": ph.placeholder_format.idx,
                "x": round(ph.left / 914400, 3) if ph.left else 0,
                "y": round(ph.top / 914400, 3) if ph.top else 0,
                "w": round(ph.width / 914400, 3) if ph.width else 0,
                "h": round(ph.height / 914400, 3) if ph.height else 0,
            })
        layouts_info.append({
            "index": i,
            "name": layout.name or f"Layout {i + 1}",
            "placeholders": ph_list,
        })

    # ── Pass 1: name-based classification (highest priority). ──
    for layout_info in layouts_info:
        role = _role_from_layout_name(layout_info["name"])
        if role and role not in layout_roles:
            layout_roles[role] = layout_info["index"]

    # ── Pass 2: placeholder-shape heuristics for anything still unassigned. ──
    for layout_info in layouts_info:
        ph_types = {p["type"] for p in layout_info["placeholders"]}
        # Recognise BOTH the python-pptx enum name and the OOXML tag name.
        has_ctr_title = bool(ph_types & {
            "CENTER_TITLE", "CTRTITLE", "ctrTitle",
        })
        has_title = has_ctr_title or bool(ph_types & {"TITLE", "title"})
        has_subtitle = bool(ph_types & {"SUBTITLE", "SUB_TITLE", "subTitle"})
        has_pic = bool(ph_types & {"PICTURE", "pic"})
        has_body = bool(ph_types & _CONTENT_PH_TYPES)
        body_count = sum(
            1 for p in layout_info["placeholders"] if p["type"] in _CONTENT_PH_TYPES
        )
        i = layout_info["index"]

        if has_ctr_title or (has_title and has_subtitle and not has_body):
            layout_roles.setdefault("title", i)
        elif has_title and not has_body and not has_subtitle and not has_pic:
            layout_roles.setdefault("section", i)
        elif has_title and body_count >= 2:
            layout_roles.setdefault("two_content", i)
            layout_roles.setdefault("content", i)
        elif has_title and has_body and has_pic:
            layout_roles.setdefault("picture", i)
            layout_roles.setdefault("two_content", i)
        elif has_title and has_body and not has_pic:
            layout_roles.setdefault("content", i)
        elif not ph_types or (
            len(ph_types) == 1 and ph_types & {"SLIDENUMBER", "sldNum"}
        ):
            layout_roles.setdefault("blank", i)

    # ── Pass 3: safety fallbacks so every role resolves to *some* layout. ──
    n_l = len(layouts_info)
    if n_l > 0:
        layout_roles.setdefault("title", 0)
        layout_roles.setdefault("content", 1 if n_l > 1 else 0)
        layout_roles.setdefault("section", layout_roles.get("title", 0))
        layout_roles.setdefault("two_content", layout_roles.get("content", 0))
        layout_roles.setdefault("blank", n_l - 1)

    # ── Colors (unchanged). ──
    try:
        import zipfile
        from xml.etree import ElementTree as ET
        with zipfile.ZipFile(template_path) as z:
            theme_files = [
                n for n in z.namelist()
                if n.startswith("ppt/theme/") and n.endswith(".xml")
            ]
            colors = dict(_DEFAULT_COLORS)
            if theme_files:
                root = ET.fromstring(z.read(theme_files[0]))
                ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
                clr = root.find(".//a:clrScheme", ns)
                if clr is not None:
                    for tag, key in [
                        ("dk1", "dark"), ("dk2", "primary"),
                        ("lt1", "accent"), ("lt2", "secondary"),
                        ("accent1", "primary"), ("accent2", "secondary"),
                    ]:
                        el = clr.find(f"a:{tag}", ns)
                        if el is not None:
                            srgb = el.find("a:srgbClr", ns)
                            if srgb is not None:
                                colors[key] = (srgb.get("val") or "").upper()
    except Exception:
        colors = dict(_DEFAULT_COLORS)

    return {
        "layouts": layouts_info,
        "layout_roles": layout_roles,
        "colors": colors,
        "template_path": template_path,
    }


async def render_pptx_from_template(
    plan: Any, template_path: str, work_dir: str,
    output_filename: str = "presentation.pptx",
) -> str | None:
    """Render the deck INTO a copy of the template. Returns path, or None to signal
    the caller to fall back to PptxGenJS."""
    import asyncio
    return await asyncio.to_thread(
        _render_sync, plan, template_path, work_dir, output_filename
    )


def _render_sync(plan, template_path, work_dir, output_filename) -> str | None:
    reuse = can_reuse_template(template_path)
    if reuse is None:
        return None

    # Prefer plan-supplied roles: engine._apply_template_theme may have
    # already refined them with an async LLM classifier. Fall back to the
    # heuristic result whenever the plan didn't override a specific role.
    plan_roles = (getattr(plan, "template_master", {}) or {}).get("layout_roles")
    if isinstance(plan_roles, dict) and plan_roles:
        reuse["layout_roles"] = {**reuse["layout_roles"], **plan_roles}

    try:
        prs = Presentation(template_path)
        # Extract per-slide libraries BEFORE stripping (backgrounds + logos).
        bg_lib = _extract_bg_library(prs)
        decor_lib = _extract_decor_library(prs)
        _strip_existing_slides(prs)
    except Exception:
        logger.debug("Failed to open template for reuse", exc_info=True)
        return None

    roles = reuse["layout_roles"]
    colors = reuse.get("colors", _DEFAULT_COLORS)
    layout_list = prs.slide_layouts

    override_bg = bool(bg_lib)
    if override_bg:
        logger.info("[template-backend] bg library: %s",
                    ", ".join(sorted(bg_lib.keys())))
    if decor_lib:
        logger.info("[template-backend] decor library: %s",
                    ", ".join(f"{k}={len(v)}" for k, v in decor_lib.items()))

    for s in plan.slides:
        role = _role_for(s.slide_type, roles)
        layout = layout_list[role]
        slide = prs.slides.add_slide(layout)

        # Background first (rendered below everything).
        if override_bg:
            entry = bg_lib.get(_bg_role_for(s.slide_type))
            applied = _apply_bg(slide, entry) if entry else False
            if not applied:
                _ensure_light_bg(slide)

        # Per-slide decorations (logos etc.). Master/layout shapes are already
        # inherited by add_slide(layout); this only adds source-slide-level pics.
        _apply_decor(slide, decor_lib.get(_bg_role_for(s.slide_type)))

        try:
            _fill_slide(slide, s, colors, plan)
        except Exception as e:
            logger.warning("[template-backend] failed to fill slide '%s': %s",
                           (s.title or "")[:40], e)

    out = os.path.join(work_dir, output_filename)
    prs.save(out)

    try:
        _sr = _load_sibling("slide_renderer")
        overflow = _sr.check_overflow_structured(out)
        if overflow:
            logger.info("[template-backend] detected %d overflows, applying fixes",
                        len(overflow))
            _apply_template_overflow_fixes(prs, plan, overflow)
            prs.save(out)
    except Exception as e:
        logger.debug("[template-backend] overflow check failed: %s", e)

    logger.info("[template-backend] rendered %d slides into template",
                len(plan.slides))
    return out

def _apply_template_overflow_fixes(prs, plan, overflow_results):
    """Post-render QA fixes for the template path.

    Because we prefer to keep the *large* font size chosen by _fit_bullets_to_box,
    we fix critical overflow by trimming content — bullet count first, then
    per-bullet length — instead of shrinking fonts below the readability floor.
    """
    for item in overflow_results:
        idx = item["slide_index"]
        if idx >= len(prs.slides) or idx >= len(plan.slides):
            continue

        severity = item.get("severity", "warning")
        overflow_in = float(item.get("overflow_inches", 0) or 0)
        if severity != "critical" and overflow_in < 0.3:
            continue

        slide = prs.slides[idx]
        slide_data = plan.slides[idx]

        # Strategy 1: aggressive bullet trim (keep 3 for critical, 4 otherwise).
        if slide_data.bullets:
            max_keep = 3 if severity == "critical" else 4
            if len(slide_data.bullets) > max_keep:
                slide_data.bullets = slide_data.bullets[:max_keep]

            # Strategy 2: shorten each bullet to ~70 chars.
            trimmed: list[str] = []
            for b in slide_data.bullets:
                if len(b) > 70:
                    cut = b[:70].rfind(" ")
                    cut = cut if cut > 35 else 70
                    trimmed.append(b[:cut].rstrip(",.;:") + "...")
                else:
                    trimmed.append(b)
            slide_data.bullets = trimmed

        # Re-fill the body placeholder; _fit_bullets_to_box will pick the
        # largest font that now fits (font size may recover after trimming).
        try:
            _fill_body_bullets(slide, slide_data)
        except Exception:
            pass

def _role_for(slide_type: str, roles: dict[str, int]) -> int:
    """Map slide_type → layout index using the classified roles."""
    mapping = {
        "title": "title",
        "section": "section",
        "conclusion": "section",  # use section-style layout for conclusion
        "content": "content",
        "content_figure": "picture",      # CHANGED: prefer picture over two_content for single figure
        "full_figure": "blank",           # blank canvas for centered figure
        "metrics": "content",             # content base, then overlay metrics
        "table": "content",               # content base, then overlay table
        "two_column": "two_content",      # ideally has two body placeholders
        "icon_rows": "content",           # content base, then overlay rows
        "steps": "content",               # content base, then overlay steps
        "emphasis_box": "content",        # content base, then overlay box
    }
    key = mapping.get(slide_type, "content")

    # Add fallback chain to avoid picking wrong layout
    fallback_chain = {
        "picture": ["picture", "two_content", "content", "blank"],
        "two_content": ["two_content", "content", "blank"],
        "content": ["content", "blank"],
        "section": ["section", "title", "blank"],
        "title": ["title", "blank"],
        "blank": ["blank", "content"],
    }

    for candidate in fallback_chain.get(key, [key, "content", "blank"]):
        if candidate in roles:
            return roles[candidate]

    # Ultimate fallback: return first available layout
    return roles.get("content", 0)


def _strip_existing_slides(prs) -> None:
    """Remove all existing slides and their parts from the template package."""
    # Iterate over a copy because we'll be mutating the list
    for sld_id_elem in list(prs.slides._sldIdLst):
        rId = sld_id_elem.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
        if rId:
            # Drop the relationship from the package, which also removes the slide part
            prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(sld_id_elem)

_KEEP_DECOR_TYPES = {"SLIDENUMBER", "DATE", "FOOTER", "slideNumber", "dt", "ftr"}

def _delete_placeholder(ph):
    """Physically remove a placeholder shape (avoids the empty-prompt render)."""
    ph._element.getparent().remove(ph._element)

def _remove_unused_placeholders(slide, filled_idxs):
    """Delete placeholders NOT in filled_idxs, except decorative ones (slide#, footer)."""
    for ph in list(slide.placeholders):
        t_str = str(ph.placeholder_format.type).split(".")[-1].split(" ")[0]
        if t_str.upper() in {x.upper() for x in _KEEP_DECOR_TYPES}:
            continue
        if ph.placeholder_format.idx not in filled_idxs:
            _delete_placeholder(ph)

def _fill_slide(slide, s, colors, plan):
    """Fill one template slide by slide_type, then delete unused placeholders."""
    filled_idxs: set[int] = set()
    stype = s.slide_type

    # ── 1. Title (works for every slide type that has a title) ──
    title_ph = _find_placeholder_by_type(
        slide, ["TITLE", "CENTER_TITLE", "CTRTITLE", "title", "ctrTitle"]
    )
    if title_ph is None and s.title:
        body_candidates = [
            ph for ph in slide.placeholders
            if str(ph.placeholder_format.type).split(".")[-1].split(" ")[0]
               in _CONTENT_PH_TYPES
        ]
        if body_candidates:
            # Pick the placeholder with the smallest ``top`` (visually the
            # top-most one) as the title fallback.
            title_ph = min(body_candidates, key=lambda p: (p.top or 0))

    if title_ph is not None and s.title:
        _fit_placeholder_text(
            title_ph, s.title, max_size=32, min_size=20,
        )
        filled_idxs.add(title_ph.placeholder_format.idx)
    elif s.title:
        # Ultimate fallback: an absolute textbox at the top of the slide.
        W, _H = _slide_dims(slide)
        tb = _new_textbox(
            slide,
            Inches(_MARGIN_L), Inches(0.35),
            Inches(W - _MARGIN_L - _MARGIN_R), Inches(0.7),
        )
        p = tb.text_frame.paragraphs[0]
        p.text = s.title
        p.font.size = Pt(_dynamic_font_size(
            s.title, W - _MARGIN_L - _MARGIN_R, 0.7,
            max_size=30, min_size=18,
        ))
        p.font.bold = True

    # ── 2. Type-specific body ──
    if stype == "title":
        # Fill subtitle if plan carries author/affiliation info.
        sub_ph = _find_placeholder_by_role(slide, "subtitle")
        sub_text = s.subtitle or (s.extra or {}).get("authors", "")
        if sub_ph is not None and sub_text:
            _fit_placeholder_text(
                sub_ph, sub_text, max_size=22, min_size=14,
            )
            filled_idxs.add(sub_ph.placeholder_format.idx)
    elif stype == "section":
        pass  # title only; other placeholders will be deleted below
    elif stype == "content":
        idx = _fill_body_bullets(slide, s)
        if idx is not None:
            filled_idxs.add(idx)
    elif stype == "content_figure":
        filled_idxs.update(_fill_content_figure(slide, s, colors))
    elif stype == "full_figure":
        _fill_full_figure(slide, s, colors)
    elif stype == "metrics":
        _fill_metrics(slide, s, colors)  # overlay only
    elif stype == "table":
        _fill_table_slide(slide, s, colors)  # overlay only
    elif stype == "two_column":
        filled_idxs.update(_fill_two_column(slide, s, colors))
    elif stype == "icon_rows":
        _fill_icon_rows(slide, s, colors)  # overlay only
    elif stype == "steps":
        _fill_steps(slide, s, colors)  # overlay only
    elif stype == "emphasis_box":
        idx = _fill_emphasis_box(slide, s, colors)
        if idx is not None:
            filled_idxs.add(idx)
    elif stype == "conclusion":
        idx = _fill_body_bullets(slide, s)
        if idx is not None:
            filled_idxs.add(idx)
        _add_thanks_line(slide)

    # Render slide-bottom citations if the planner attached any. These are
    # absolute textboxes; they do not consume a placeholder idx.
    if s.slide_type not in ("title", "references") and s.citations:
        _add_slide_footnotes(slide, s.citations, colors)

    # Dedicated References slide handler.
    if s.slide_type == "references":
        _fill_references_slide(slide, s, plan.references or [], colors)

    # ── 3. Remove every placeholder we did NOT fill (keeps decor like page#) ──
    _remove_unused_placeholders(slide, filled_idxs)



# ── Placeholder helpers ──

def _find_placeholder_by_type(slide, type_names) -> Any | None:
    """Find the first placeholder whose type name matches one of the given names."""
    type_set = {t.upper() for t in type_names}
    for ph in slide.placeholders:
        ph_type_str = str(ph.placeholder_format.type).split(".")[-1].split(" ")[0].upper()
        # Also check exact match from idx-based mapping
        if ph_type_str in type_set:
            return ph
    return None


def _find_placeholder_by_role(slide, role: str) -> Any | None:
    """Find a placeholder by semantic role (title, body, subtitle, picture)."""
    type_map = {
        # python-pptx stringifies PP_PLACEHOLDER as e.g. "CENTER_TITLE (15)",
        # NOT the OOXML tag "ctrTitle". Include both forms so this helper works
        # whether callers came in through python-pptx enums or raw XML.
        "title": ["TITLE", "CENTER_TITLE", "CTRTITLE", "ctrTitle"],
        # Templates commonly use OBJECT ("Content Placeholder") for body text.
        "body": ["BODY", "OBJECT", "body", "object"],
        "subtitle": ["SUBTITLE", "SUB_TITLE", "subTitle"],
        "picture": ["PICTURE", "pic"],
    }
    return _find_placeholder_by_type(slide, type_map.get(role, []))


def _fill_body_bullets(
    slide, s,
    candidates: tuple[int, ...] = (24, 22, 20, 18),
) -> int | None:
    """Fill the body placeholder with bullets at the largest font from
    *candidates* that fits. If the floor still overflows, trim trailing
    bullets in-place (mutating s.bullets) so downstream QA sees the actual
    rendered content.

    Callers that render inside a narrower column (e.g. content_figure with a
    side picture) can pass smaller *candidates* like (20, 18, 16).
    """
    if not s.bullets:
        return None

    body_ph = _find_placeholder_by_role(slide, "body")
    if body_ph is not None and body_ph.has_text_frame:
        tf = body_ph.text_frame
        tf.clear()
        tf.word_wrap = True
        try:
            tf.vertical_anchor = MSO_ANCHOR.TOP
        except Exception:
            pass

        avail_h = (body_ph.height / 914400) if body_ph.height else 3.5
        box_w = (body_ph.width / 914400) if body_ph.width else 8.0

        fitted, fs, para = _fit_bullets_to_box(
            s.bullets, avail_h, box_w, candidates=candidates,
        )
        s.bullets = fitted  # persist trim so QA overflow check is coherent

        for i, b in enumerate(s.bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = b
            p.font.size = Pt(fs)
            p.font.color.rgb = _hex_to_rgb("2D2D2D")
            p.space_after = Pt(para)
            p.level = 0
        return body_ph.placeholder_format.idx

    # No body placeholder: fall back to an absolute textbox sized from real slide.
    _add_bullets_absolute(slide, s.bullets)
    return None

def _fit_bullets_to_box(
    bullets: list[str],
    avail_h: float,
    box_w: float,
    candidates: tuple[int, ...] = (20, 18, 16),
) -> tuple[list[str], int, int]:
    """Pick the largest font from *candidates* that lets *bullets* fit in the
    box. If even the floor overflows, trim trailing bullets (keeping ≥1)
    rather than shrink below the floor — a template deck is read at
    projection distance and 18pt is our readability floor.

    Returns: (possibly-trimmed bullets, chosen font size, paragraph spacing pt).
    """
    bullets = list(bullets)
    fill_target = 0.92  # aim to use ≤92% of vertical space

    for fs in candidates:
        est = _estimate_bullet_height(bullets, fs, box_w)
        if est <= avail_h * fill_target:
            para = _distribute_para_space(len(bullets), est, avail_h)
            return bullets, fs, para

    # Floor reached; trim from the end until content fits at the floor size.
    floor = candidates[-1]
    while len(bullets) > 1:
        bullets.pop()
        if _estimate_bullet_height(bullets, floor, box_w) <= avail_h * fill_target:
            break
    est = _estimate_bullet_height(bullets, floor, box_w)
    para = _distribute_para_space(len(bullets), est, avail_h)
    return bullets, floor, para


def _distribute_para_space(n_bullets: int, used_h: float, avail_h: float) -> int:
    """Distribute leftover vertical space as paragraph spacing (bounded)."""
    if n_bullets <= 1:
        return 8
    remaining_in = max(0.0, avail_h - used_h)
    # Convert to points, share across gaps, cap at 28pt to avoid loose looks.
    return min(28, max(6, int(remaining_in * 72 / (n_bullets - 1) * 0.4)))

def _estimate_bullet_height(bullets, font_size_pt, box_width_in):
    """Estimate rendered bullet-list height in inches.

    We ceil() line counts and add a 15% word-wrap safety margin so the fitting
    loop never *undershoots* the real height. Undershoot means overflow at
    render time — which we cannot fix visually because there is no image
    critique. Overshoot only costs us a font-size step, which is acceptable.
    """
    total_lines = 0
    for b in bullets:
        # Rough width per char: full-width for CJK, ~0.55 em for Latin.
        char_count = sum(1.0 if ord(c) >= 0x4E00 else 0.55 for c in b)
        raw = char_count * font_size_pt / max(0.5, (box_width_in * 72)) * 1.15
        total_lines += max(1, math.ceil(raw))
    # Line height 1.4 × font size, plus a small internal padding allowance.
    return (total_lines * font_size_pt * 1.4) / 72 + 0.1



# ── Complex layout fillers ──

def _fill_content_figure(slide, s, colors) -> set[int]:
    filled: set[int] = set()
    body_ph = _find_placeholder_by_role(slide, "body")
    pic_ph = _find_placeholder_by_role(slide, "picture")

    has_fig_file = bool(s.figure_path) and os.path.exists(s.figure_path)

    if pic_ph is not None and has_fig_file:
        # Template already has a picture placeholder — the layout defines
        # the geometry, so bullets vs picture never overlap.
        try:
            pic_ph.insert_picture(s.figure_path)
            filled.add(pic_ph.placeholder_format.idx)
        except Exception:
            _add_picture_absolute(slide, s.figure_path, side="right")
        if body_ph is not None:
            idx = _fill_body_bullets(slide, s)
            if idx is not None:
                filled.add(idx)
    else:
        # No picture placeholder available: shrink the body placeholder to
        # the left half BEFORE filling bullets, then drop the picture on
        # the right half. Without the shrink, bullets stretch full-width
        # and collide with the absolute-positioned picture.
        if has_fig_file and body_ph is not None and body_ph.width is not None:
            W, _H = _slide_dims(slide)
            # Content region right edge in EMU (1 in = 914400 EMU).
            content_right_emu = int((W - _MARGIN_R) * 914400)
            # Materialise every geometry field FIRST so the placeholder no
            # longer relies on layout inheritance. Writing .width alone drops
            # the other three from the new xfrm and the placeholder collapses
            # to height=0 (same bug class as _fill_emphasis_box).
            left0 = body_ph.left or int(_MARGIN_L * 914400)
            top0 = body_ph.top
            width0 = body_ph.width
            height0 = body_ph.height
            # Keep bullets to ~48% of the content width so there is room
            # for the picture + a small gutter on the right.
            max_body_w = int((W - _MARGIN_L - _MARGIN_R) * 0.48 * 914400)
            new_w = min(width0, max_body_w)
            # Safety: never overflow the content region.
            if left0 + new_w > content_right_emu:
                new_w = max(1, content_right_emu - left0)
            # Assign all four dims so spPr/xfrm carries off + ext together.
            body_ph.left = left0
            if top0 is not None:
                body_ph.top = top0
            body_ph.width = new_w
            if height0 is not None:
                body_ph.height = height0

        idx = _fill_body_bullets(slide, s)
        if idx is not None:
            filled.add(idx)
        if has_fig_file:
            _add_picture_absolute(slide, s.figure_path, side="right")

    if s.figure_caption:
        _add_caption(slide, s.figure_caption)
    return filled


def _fill_full_figure(slide, s, colors):
    """Full-width figure slide."""
    if s.figure_path and os.path.exists(s.figure_path):
        _add_picture_absolute(slide, s.figure_path, side="center")
    if s.figure_caption:
        _add_caption(slide, s.figure_caption)


def _fill_metrics(slide, s, colors):
    """Metrics: overlay metric cards on top of content layout."""
    _fill_body_bullets(slide, s)  # title already set
    if s.metrics:
        _add_metrics_absolute(slide, s, colors)


def _fill_table_slide(slide, s, colors):
    """Table slide: overlay table shape."""
    if s.table_headers and s.table_rows:
        _add_table_absolute(slide, s, colors)


def _fill_two_column(slide, s, colors) -> set[int]:
    """Two-column slide.

    Path A: use N leftmost body/object placeholders when the template layout
            provides them.
    Path B: fall back to absolute textboxes sized from the real slide.
    Font sizes are chosen dynamically per column based on the actual box width.
    """
    columns = (s.extra or {}).get("columns", [])
    filled: set[int] = set()
    if not columns:
        return filled

    n = len(columns)

    # ── Path A: template placeholders ──
    body_phs = [
        ph for ph in slide.placeholders
        if str(ph.placeholder_format.type).split(".")[-1].split(" ")[0]
        in _CONTENT_PH_TYPES
    ]
    body_phs.sort(key=lambda p: (p.left or 0))

    if len(body_phs) >= n:
        for ph, col in zip(body_phs[:n], columns, strict=True):
            tf = ph.text_frame
            tf.clear()
            tf.word_wrap = True
            try:
                tf.vertical_anchor = MSO_ANCHOR.TOP
            except Exception:
                pass

            col_w_in = (ph.width / 914400) if ph.width else 4.0
            bullets = col.get("bullets") or []

            # One uniform font per column (based on the longest bullet).
            bullet_fs = _uniform_font_size(
                bullets, col_w_in, max_size=18, min_size=13,
            )

            if col.get("sub_title"):
                p = tf.paragraphs[0]
                p.text = col["sub_title"]
                p.font.bold = True
                p.font.size = Pt(_dynamic_font_size(
                    col["sub_title"], col_w_in,
                    max_size=20, min_size=14,
                ))
                p.font.color.rgb = _hex_to_rgb(colors["primary"])
                start = 1
            else:
                start = 0

            for i, b in enumerate(bullets):
                p = tf.paragraphs[0] if (i + start) == 0 else tf.add_paragraph()
                p.text = f"• {b}"
                p.font.size = Pt(bullet_fs)
                p.font.color.rgb = _hex_to_rgb(colors["bodyText"])

            if col.get("figure_path") and os.path.exists(col["figure_path"]):
                try:
                    fx = ph.left / 914400
                    fy = (ph.top + ph.height) / 914400 - 1.6
                    fw = ph.width / 914400
                    slide.shapes.add_picture(
                        col["figure_path"],
                        Inches(fx), Inches(fy), width=Inches(fw),
                    )
                except Exception:
                    pass

            filled.add(ph.placeholder_format.idx)

        note = (s.extra or {}).get("emphasis_note")
        if note:
            W, H = _slide_dims(slide)
            tb = _new_textbox(
                slide,
                Inches(_MARGIN_L), Inches(H - 0.6),
                Inches(W - _MARGIN_L - _MARGIN_R), Inches(0.4),
            )
            p = tb.text_frame.paragraphs[0]
            p.text = note
            p.font.size = Pt(_dynamic_font_size(
                note, W - _MARGIN_L - _MARGIN_R,
                max_size=13, min_size=10,
            ))
            p.font.italic = True
            p.font.color.rgb = _hex_to_rgb(colors["muted"])
        return filled

    # ── Path B: absolute textboxes ──
    W, H = _slide_dims(slide)
    cx, cy, cw, _ch = _content_region(slide)
    gutter = 0.25
    col_w = (cw - gutter * (n - 1)) / n

    for i, col in enumerate(columns):
        x = cx + i * (col_w + gutter)
        y = cy
        bullets = col.get("bullets") or []
        bullet_fs = _uniform_font_size(
            bullets, col_w, max_size=17, min_size=13,
        )

        if col.get("sub_title"):
            tb = _new_textbox(
                slide,
                Inches(x), Inches(y), Inches(col_w), Inches(0.4),
            )
            p = tb.text_frame.paragraphs[0]
            p.text = col["sub_title"]
            p.font.size = Pt(_dynamic_font_size(
                col["sub_title"], col_w,
                max_size=20, min_size=14,
            ))
            p.font.bold = True
            p.font.color.rgb = _hex_to_rgb(colors["primary"])
            y += 0.5

        for b in bullets:
            tb = _new_textbox(
                slide,
                Inches(x + 0.1), Inches(y), Inches(col_w - 0.1), Inches(0.35),
            )
            p = tb.text_frame.paragraphs[0]
            p.text = f"• {b}"
            p.font.size = Pt(bullet_fs)
            p.font.color.rgb = _hex_to_rgb(colors["bodyText"])
            y += 0.35

        if col.get("figure_path") and os.path.exists(col["figure_path"]):
            try:
                slide.shapes.add_picture(
                    col["figure_path"],
                    Inches(x), Inches(y), width=Inches(col_w),
                )
            except Exception:
                pass

    note = (s.extra or {}).get("emphasis_note")
    if note:
        tb = _new_textbox(
            slide,
            Inches(cx), Inches(H - 0.6), Inches(cw), Inches(0.4),
        )
        p = tb.text_frame.paragraphs[0]
        p.text = note
        p.font.size = Pt(_dynamic_font_size(
            note, cw, max_size=13, min_size=10,
        ))
        p.font.italic = True
        p.font.color.rgb = _hex_to_rgb(colors["muted"])
    return filled


def _fill_icon_rows(slide, s, colors):
    rows = (s.extra or {}).get("rows", [])
    if not rows:
        return

    cx, cy, cw, ch = _content_region(slide)
    label_w = 0.85
    label_gap = 0.25
    text_x = cx + label_w + label_gap
    text_w = cw - label_w - label_gap

    row_h = 1.15
    row_gap = 0.15
    total_h = len(rows) * row_h + max(0, len(rows) - 1) * row_gap
    if total_h > ch:
        scale = ch / total_h
        row_h *= scale
        row_gap *= scale
        total_h = ch

    start_y = cy + max(0.0, (ch - total_h) / 2)

    # Uniform font sizes across all rows for a consistent look.
    label_fs = _uniform_font_size(
        [r.get("label", "") for r in rows], label_w,
        max_size=22, min_size=14,
    )
    header_fs = _uniform_font_size(
        [r.get("header", "") for r in rows], text_w,
        max_size=18, min_size=13,
    )
    desc_fs = _uniform_font_size(
        [r.get("description", "") for r in rows], text_w,
        box_height_in=row_h - 0.5,
        max_size=13, min_size=10,
    )

    for i, row in enumerate(rows):
        y = start_y + i * (row_h + row_gap)

        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(cx), Inches(y), Inches(label_w), Inches(row_h),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = _hex_to_rgb(colors["tableFill"])
        shape.line.fill.background()
        shape.text_frame.text = row.get("label", str(i + 1))
        shape.text_frame.paragraphs[0].font.size = Pt(label_fs)
        shape.text_frame.paragraphs[0].font.bold = True
        shape.text_frame.paragraphs[0].font.color.rgb = _hex_to_rgb(colors["primary"])
        shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        if row.get("header"):
            tb = _new_textbox(
                slide,
                Inches(text_x), Inches(y + 0.05), Inches(text_w), Inches(0.4),
            )
            p = tb.text_frame.paragraphs[0]
            p.text = row["header"]
            p.font.size = Pt(header_fs)
            p.font.bold = True
            p.font.color.rgb = _hex_to_rgb(colors["primary"])

        if row.get("description"):
            tb = _new_textbox(
                slide,
                Inches(text_x), Inches(y + 0.48), Inches(text_w), Inches(0.55),
            )
            p = tb.text_frame.paragraphs[0]
            p.text = row["description"]
            p.font.size = Pt(desc_fs)
            p.font.color.rgb = _hex_to_rgb(colors["bodyText"])

def _fill_steps(slide, s, colors):
    steps = (s.extra or {}).get("steps", [])
    if not steps:
        return

    cx, cy, cw, ch = _content_region(slide)
    n = len(steps)
    step_gap = 0.25
    step_w = (cw - step_gap * (n - 1)) / n
    card_h = min(2.7, ch * 0.85)
    card_y = cy + max(0.0, (ch - card_h) / 2)

    # Uniform sizes per element type across all step cards.
    num_fs = _uniform_font_size(
        [st.get("step_number", "") for st in steps], step_w * 0.8,
        max_size=30, min_size=18,
    )
    title_fs = _uniform_font_size(
        [st.get("step_title", "") for st in steps], step_w - 0.3,
        box_height_in=0.5,
        max_size=20, min_size=13,
    )
    desc_fs = _uniform_font_size(
        [st.get("step_desc", "") for st in steps], step_w - 0.3,
        box_height_in=max(0.4, card_h - 1.3),
        max_size=13, min_size=9,
    )

    for i, step in enumerate(steps):
        x = cx + i * (step_w + step_gap)

        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(card_y), Inches(step_w), Inches(card_h),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = _hex_to_rgb(
            colors["primary"] if i == 0 else colors["tableFill"]
        )
        shape.line.fill.background()

        is_dark = (i == 0)
        text_color = _hex_to_rgb("FFFFFF" if is_dark else colors["primary"])
        muted_color = _hex_to_rgb(colors["secondary"] if is_dark else colors["bodyText"])

        tb = _new_textbox(
            slide,
            Inches(x + 0.1), Inches(card_y + 0.15),
            Inches(step_w - 0.2), Inches(0.4),
        )
        p = tb.text_frame.paragraphs[0]
        p.text = step.get("step_number", str(i + 1))
        p.font.size = Pt(num_fs)
        p.font.bold = True
        p.font.color.rgb = text_color
        p.alignment = PP_ALIGN.CENTER

        if step.get("step_title"):
            tb = _new_textbox(
                slide,
                Inches(x + 0.15), Inches(card_y + 0.6),
                Inches(step_w - 0.3), Inches(0.55),
            )
            p = tb.text_frame.paragraphs[0]
            p.text = step["step_title"]
            p.font.size = Pt(title_fs)
            p.font.bold = True
            p.font.color.rgb = text_color
            p.alignment = PP_ALIGN.CENTER

        if step.get("step_desc"):
            tb = _new_textbox(
                slide,
                Inches(x + 0.15), Inches(card_y + 1.2),
                Inches(step_w - 0.3), Inches(max(0.4, card_h - 1.3)),
            )
            p = tb.text_frame.paragraphs[0]
            p.text = step["step_desc"]
            p.font.size = Pt(desc_fs)
            p.font.color.rgb = muted_color
            p.alignment = PP_ALIGN.CENTER

        if i < n - 1:
            tb = _new_textbox(
                slide,
                Inches(x + step_w + step_gap * 0.3),
                Inches(card_y + card_h / 2 - 0.2),
                Inches(step_gap * 0.4), Inches(0.4),
            )
            p = tb.text_frame.paragraphs[0]
            p.text = "→"
            p.font.size = Pt(18)
            p.font.color.rgb = _hex_to_rgb(colors["muted"])
            p.alignment = PP_ALIGN.CENTER

def _fill_emphasis_box(slide, s, colors) -> int | None:
    """Emphasis-box slide: bullets on top, prominent colored box at the bottom.

    Reserves vertical space for the box FIRST so bullets don't overlap it on
    tall (16:9 wide) slides. Returns the filled body-placeholder idx (or None
    if a bullet-placeholder was not used).
    """
    W, H = _slide_dims(slide)
    box_h = 1.15
    box_reserve = box_h + 0.25  # box height + gap above it

    # ── 1. Fill bullets (into the body placeholder if the layout provides one) ──
    # We rely on _fill_body_bullets' own logic to detect body/object placeholder
    # and to top-anchor its text_frame; it returns the filled placeholder idx.
    body_ph = _find_placeholder_by_role(slide, "body")
    # Snapshot every geometry field BEFORE mutating. python-pptx resolves
    # placeholder inheritance from the layout lazily, so writing only
    # .height causes cx / off/x / off/y to be dropped from the new xfrm and
    # the placeholder renders at width=0 in PowerPoint / WPS.
    saved_geom = None
    if body_ph is not None and body_ph.height is not None:
        left0 = body_ph.left
        top0 = body_ph.top
        width0 = body_ph.width
        height0 = body_ph.height
        if None not in (left0, top0, width0, height0):
            saved_geom = (left0, top0, width0, height0)
            try:
                new_h = int(height0 - int(box_reserve * 914400))
                # Write ALL four dims together so the placeholder's spPr
                # ends up with a complete <a:xfrm><a:off/><a:ext/></a:xfrm>
                # instead of a lone <a:ext cy=".."/>.
                body_ph.left = left0
                body_ph.top = top0
                body_ph.width = width0
                body_ph.height = new_h
            except Exception:
                saved_geom = None
    filled_idx = _fill_body_bullets(slide, s)
    if saved_geom is not None:
        try:
            body_ph.left, body_ph.top, body_ph.width, body_ph.height = saved_geom
        except Exception:
            pass

    box_text = (s.extra or {}).get("box_text")
    if not box_text:
        return filled_idx

    # ── 2. Emphasis box pinned to the real slide bottom ──
    box_x = _MARGIN_L - 0.1
    box_w = W - _MARGIN_L - _MARGIN_R + 0.2
    box_y = H - box_reserve

    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(box_x), Inches(box_y), Inches(box_w), Inches(box_h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _hex_to_rgb(colors["tableFill"])
    shape.line.fill.background()

    # Accent bar on the left edge, same height as the box.
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(box_x), Inches(box_y), Inches(0.06), Inches(box_h),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = _hex_to_rgb(colors["primary"])
    bar.line.fill.background()

    # Text inside the box, indented past the accent bar.
    tb = _new_textbox(
        slide,
        Inches(box_x + 0.25), Inches(box_y + 0.1),
        Inches(box_w - 0.35), Inches(box_h - 0.2),
    )
    tf = tb.text_frame
    try:
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.text = box_text
    p.font.size = Pt(_dynamic_font_size(
        box_text, box_w - 0.35,
        box_height_in=box_h - 0.2,
        max_size=18, min_size=12,
    ))
    p.font.bold = True
    p.font.color.rgb = _hex_to_rgb(colors["primary"])
    return filled_idx



# ── Absolute-position helpers ──

def _add_bullets_absolute(slide, bullets, font_size=18):
    """Add bullets as absolute textbox."""
    if not bullets:
        return
    W, H = _slide_dims(slide)
    x = _MARGIN_L
    y = _BODY_TOP
    w = W - _MARGIN_L - _MARGIN_R
    h = H - _BODY_TOP - _BOTTOM
    tb = _new_textbox(slide, Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {b}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = _hex_to_rgb("2D2D2D")


def _add_picture_absolute(slide, path, side="right"):
    """Place a figure preserving aspect ratio, dimensioned from the real slide size."""
    from PIL import Image
    try:
        with Image.open(path) as im:
            w, h = im.size
        ratio = w / h if h else 1.6
    except Exception:
        ratio = 1.6

    # Derive picture box from the actual slide dimensions.
    W, H = _slide_dims(slide)
    cx, cy, cw, ch = _content_region(slide)

    if side == "center":
        # Center pic uses full content width, with room for a caption below.
        max_w = cw
        max_h = ch - 0.5           # leave space for a caption
        fw = min(max_w, max_h * ratio)
        fh = fw / ratio
        x = (W - fw) / 2
        y = cy + (max_h - fh) / 2
    else:
        # Right-side pic takes ~half the content width.
        max_w = cw * 0.5
        max_h = ch - 0.4
        fw = min(max_w, max_h * ratio)
        fh = fw / ratio
        x = W - _MARGIN_R - fw
        y = cy + (max_h - fh) / 2

    slide.shapes.add_picture(
        path, Inches(max(0.3, x)), Inches(max(cy, y)),
        width=Inches(fw), height=Inches(fh),
    )


def _add_caption(slide, text):
    """Add small italic caption anchored to the slide's actual bottom."""
    if not text:
        return
    # 5.05" bottom-line was correct for a 5.625" slide, but on a 7.5" wide-screen
    # template it lands in the vertical middle. Anchor to the real slide height.
    W, H = _slide_dims(slide)
    tb = _new_textbox(
        slide,
        Inches(_MARGIN_L), Inches(H - 0.55),
        Inches(W - _MARGIN_L - _MARGIN_R), Inches(0.35),
    )
    p = tb.text_frame.paragraphs[0]
    p.text = text[:160]
    p.font.size = Pt(_dynamic_font_size(
        text[:160], W - _MARGIN_L - _MARGIN_R,
        box_height_in=0.35, max_size=13, min_size=10,
    ))
    p.font.italic = True
    p.font.color.rgb = _hex_to_rgb("6B7280")

def _add_metrics_absolute(slide, s, colors):
    """Overlay metric cards, sized from the real slide width. Value/label
    font sizes are picked once per slide based on the longest string so the
    cards look uniform even when one metric is much longer than the others."""
    metrics = s.metrics or []
    n = max(1, len(metrics))
    cx, cy, cw, ch = _content_region(slide)
    card_gap = 0.3
    card_w = (cw - card_gap * (n - 1)) / n
    card_h = 1.75
    card_y = cy + max(0.0, (ch - card_h) / 2)

    value_fs = _uniform_font_size(
        [str(m.get("value", "")) for m in metrics], card_w - 0.1,
        max_size=44, min_size=24,
    )
    label_fs = _uniform_font_size(
        [str(m.get("label", "")) for m in metrics], card_w - 0.1,
        max_size=16, min_size=11,
    )

    for i, m in enumerate(metrics):
        x = cx + i * (card_w + card_gap)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(card_y), Inches(card_w), Inches(card_h),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = _hex_to_rgb(colors["tableFill"])
        shape.line.fill.background()

        tb_val = _new_textbox(
            slide,
            Inches(x), Inches(card_y + 0.2), Inches(card_w), Inches(0.9),
        )
        p_val = tb_val.text_frame.paragraphs[0]
        p_val.text = m.get("value", "")
        p_val.font.size = Pt(value_fs)
        p_val.font.bold = True
        p_val.font.color.rgb = _hex_to_rgb(colors["primary"])
        p_val.alignment = PP_ALIGN.CENTER

        tb_lbl = _new_textbox(
            slide,
            Inches(x), Inches(card_y + 1.15), Inches(card_w), Inches(0.5),
        )
        p_lbl = tb_lbl.text_frame.paragraphs[0]
        p_lbl.text = m.get("label", "")
        p_lbl.font.size = Pt(label_fs)
        p_lbl.font.color.rgb = _hex_to_rgb(colors["muted"])
        p_lbl.alignment = PP_ALIGN.CENTER

def _add_table_absolute(slide, s, colors):
    """Overlay a table shape, dimensioned from the real slide size."""
    headers = s.table_headers or []
    rows = (s.table_rows or [])[:7]
    cols = len(headers)
    n_rows = len(rows) + 1
    if n_rows < 2 or cols < 1:
        return

    # Content region gives us the real available width/height.
    cx, cy, cw, ch = _content_region(slide)
    row_h = min(0.4, (ch - 0.3) / n_rows)
    total_h = row_h * n_rows
    # Center vertically inside the content region when there is slack.
    table_y = cy + max(0.0, (ch - total_h) * 0.3)

    tbl_shape = slide.shapes.add_table(
        n_rows, cols,
        Inches(cx), Inches(table_y), Inches(cw), Inches(total_h),
    )
    table = tbl_shape.table

    # Header row
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = str(h)
        _style_cell(cell, bold=True, fill=colors["tableHead"], font_color="FFFFFF")

    # Data rows
    for r, row in enumerate(rows):
        is_hl = (r == s.highlight_row)
        for c, val in enumerate(row[:cols]):
            cell = table.cell(r + 1, c)
            cell.text = str(val)
            _style_cell(cell, bold=is_hl, fill=colors.get("tableFill", "F0F4FF") if r % 2 == 0 else "FFFFFF",
                        font_color=colors["primary"] if is_hl else colors["bodyText"])


def _style_cell(cell, bold=False, fill="FFFFFF", font_color="2D2D2D", size=12):
    """Apply basic styling to a table cell."""
    cell.fill.solid()
    cell.fill.fore_color.rgb = _hex_to_rgb(fill)
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = _hex_to_rgb(font_color)


def _add_thanks_line(slide):
    """Add 'Thank you' line pinned to the actual slide bottom."""
    W, H = _slide_dims(slide)
    tb = _new_textbox(
        slide,
        Inches(_MARGIN_L), Inches(H - 0.7),
        Inches(W - _MARGIN_L - _MARGIN_R), Inches(0.4),
    )
    p = tb.text_frame.paragraphs[0]
    p.text = "Thank you · Questions?"
    p.font.size = Pt(16)
    p.font.italic = True
    p.font.color.rgb = _hex_to_rgb("6B7280")

# ── Per-slide background library ─────────────────────────────────────

def _extract_bg_library(src_prs) -> dict[str, dict]:
    """Extract per-slide backgrounds keyed by role.

    For each source slide we try, in order:
      1. Slide-level <p:bg> element (solid / gradient / blipFill).
      2. A full-slide picture SHAPE (>= 80% area, positioned near the
         origin) — a lot of WPS/PowerPoint templates use a picture shape
         as the visual background instead of a real <p:bg>, and the
         previous version silently dropped those.

    Each entry carries a "kind" tag so _apply_bg knows which reattach
    strategy to run.
    """
    slide_w = src_prs.slide_width or 1
    slide_h = src_prs.slide_height or 1
    slide_area = slide_w * slide_h

    def _grab_bg_element(slide) -> dict | None:
        cSld = slide._element.find(qn("p:cSld"))
        if cSld is None:
            return None
        bg = cSld.find(qn("p:bg"))
        if bg is None:
            return None
        entry = {
            "kind": "bg_element",
            "xml": deepcopy(bg),
            "image_blob": None,
            "content_type": None,
        }
        blip = bg.find(".//" + qn("a:blip"))
        if blip is not None:
            rid = blip.get(qn("r:embed"))
            if rid:
                try:
                    part = slide.part.rels[rid].target_part
                    entry["image_blob"] = part.blob
                    entry["content_type"] = getattr(part, "content_type", None)
                except (KeyError, AttributeError):
                    pass
        return entry

    def _grab_full_slide_picture(slide) -> dict | None:
        """Largest full-slide picture shape, if any qualifies as background."""
        best: dict | None = None
        best_area = 0
        for shape in slide.shapes:
            if shape.is_placeholder:
                continue
            try:
                blob = shape.image.blob
                ct = shape.image.content_type or ""
            except (AttributeError, Exception):
                continue
            if ct.startswith("image/svg"):
                continue
            if not shape.width or not shape.height:
                continue
            area = shape.width * shape.height
            # Require ≥80% coverage AND near-origin placement so we don't
            # accidentally pick up a centred hero graphic.
            near_origin = (
                (shape.left or 0) < slide_w * 0.10
                and (shape.top or 0) < slide_h * 0.10
            )
            if area >= slide_area * 0.80 and near_origin and area > best_area:
                best = {
                    "kind": "picture_shape",
                    "image_blob": blob,
                    "content_type": ct,
                }
                best_area = area
        return best

    def _grab(slide) -> dict | None:
        entry = _grab_bg_element(slide)
        if entry is not None:
            return entry
        return _grab_full_slide_picture(slide)

    lib: dict[str, dict] = {}
    slides = list(src_prs.slides)
    if not slides:
        return lib

    cover = _grab(slides[0])
    if cover:
        lib["cover"] = cover

    if len(slides) > 1:
        closing = _grab(slides[-1])
        if closing:
            lib["closing"] = closing

    # "content" prefers a non-image bg (solid/gradient) so we don't stamp a
    # cover photo behind every body page. Fall back to whatever the cover
    # gave us if nothing else exists.
    for s in slides:
        e = _grab_bg_element(s)
        if e and e.get("image_blob") is None:
            lib["content"] = e
            break

    lib.setdefault("content", lib.get("cover"))
    lib.setdefault("section", lib.get("content"))
    return {k: v for k, v in lib.items() if v is not None}


def _add_image_rel(slide, blob: bytes, content_type: str | None) -> str:
    """Import *blob* into the slide's package and return a new relationship id.

    Trick: add_picture registers both the image part and the slide→image
    relationship. We keep the relationship (needed by the <a:blip> in the
    background XML) but remove the picture shape itself. The image bytes
    stay in the package as an unreferenced image part reachable only
    through the new rId, which is exactly what we want.
    """
    import os as _os
    import tempfile

    ext = "png"
    if content_type and "/" in content_type:
        ext = content_type.split("/", 1)[1] or "png"
    fd, path = tempfile.mkstemp(suffix=f".{ext}")
    _os.close(fd)
    with open(path, "wb") as f:
        f.write(blob)
    try:
        pic = slide.shapes.add_picture(path, 0, 0, width=1, height=1)
        blip_el = pic._element.find(".//" + qn("a:blip"))
        rid = blip_el.get(qn("r:embed")) if blip_el is not None else ""
        # Drop the placeholder shape; keep the image relationship + part.
        pic._element.getparent().remove(pic._element)
        return rid
    finally:
        try:
            _os.remove(path)
        except OSError:
            pass


def _apply_bg(slide, entry: dict | None) -> bool:
    """Attach the appropriate background to *slide* based on entry kind."""
    if not entry:
        return False
    kind = entry.get("kind", "bg_element")
    if kind == "picture_shape":
        return _apply_picture_shape_bg(slide, entry)
    return _apply_bg_element(slide, entry)


def _apply_bg_element(slide, entry: dict) -> bool:
    """Re-attach a <p:bg> XML fragment to a fresh slide.

    Image blobs (from blipFill backgrounds) are re-imported into the
    destination slide's package so the r:embed relationship resolves in
    the new scope.
    """
    cSld = slide._element.find(qn("p:cSld"))
    if cSld is None or entry.get("xml") is None:
        return False

    for old in cSld.findall(qn("p:bg")):
        cSld.remove(old)

    bg = deepcopy(entry["xml"])
    blip = bg.find(".//" + qn("a:blip"))
    if blip is not None and entry.get("image_blob"):
        try:
            new_rid = _add_image_rel(
                slide, entry["image_blob"], entry.get("content_type"),
            )
        except Exception:
            new_rid = ""
        if new_rid:
            blip.set(qn("r:embed"), new_rid)
        else:
            return False  # blip would dangle; caller falls back to light bg

    cSld.insert(0, bg)
    return True


def _apply_picture_shape_bg(slide, entry: dict) -> bool:
    """Re-render a full-slide picture as the slide's visual background.

    We drop it in as a normal picture shape sized to the slide, then push
    it to the very back of the z-order so all placeholder text renders
    on top of it. This mirrors what many WPS templates actually do.
    """
    import os as _os
    import tempfile

    blob = entry.get("image_blob")
    if not blob:
        return False
    ct = entry.get("content_type") or ""
    ext = ct.split("/", 1)[1] if "/" in ct else "png"
    fd, path = tempfile.mkstemp(suffix=f".{ext}")
    _os.close(fd)
    try:
        with open(path, "wb") as f:
            f.write(blob)
        W, H = _slide_dims(slide)
        pic = slide.shapes.add_picture(
            path, 0, 0, width=Inches(W), height=Inches(H),
        )
        _send_shape_to_back(pic)
        return True
    except Exception:
        return False
    finally:
        try:
            _os.remove(path)
        except OSError:
            pass


def _ensure_light_bg(slide) -> None:
    """Force a white <p:bg> on the slide so it can never inherit a dark master.

    Used as fallback when the background library does not carry a matching
    role. This is the minimum viable safety net referenced in the analysis
    doc ('Step 1').
    """
    cSld = slide._element.find(qn("p:cSld"))
    if cSld is None:
        return
    for old in cSld.findall(qn("p:bg")):
        cSld.remove(old)
    bg_xml = (
        '<p:bg xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
        'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        '<p:bgPr><a:solidFill><a:srgbClr val="FFFFFF"/></a:solidFill>'
        '<a:effectLst/></p:bgPr></p:bg>'
    )
    cSld.insert(0, etree.fromstring(bg_xml))


def _bg_role_for(slide_type: str) -> str:
    """Map a slide_type to a background role."""
    if slide_type == "title":
        return "cover"
    if slide_type == "conclusion":
        return "closing"
    if slide_type == "section":
        return "section"
    return "content"

# ── Dynamic font sizing ──────────────────────────────────────────────

def _dynamic_font_size(
    text: str,
    box_width_in: float,
    box_height_in: float | None = None,
    *,
    max_size: int = 24,
    min_size: int = 12,
    line_height_factor: float = 1.35,
) -> int:
    """Pick the largest integer font size (pt) at which *text* fits in the box.

    Width estimation: ~1.0em per CJK char, ~0.55em per Latin char.
    If *box_height_in* is None we only require the text to fit on ONE line
    (useful for headers/labels). If given, we allow wrapping and require the
    resulting stacked height to fit.
    Falls back to *min_size* when nothing fits.
    """
    text = (text or "").strip()
    if not text:
        return max_size

    char_units = sum((1.0 if ord(c) >= 0x4E00 else 0.55) for c in text)
    usable_w_pt = max(1.0, box_width_in * 72)

    for fs in range(max_size, min_size - 1, -1):
        w_pt = char_units * fs
        if box_height_in is None:
            if w_pt <= usable_w_pt:
                return fs
        else:
            lines = max(1, math.ceil(w_pt / usable_w_pt))
            needed_h_pt = lines * fs * line_height_factor
            if needed_h_pt <= box_height_in * 72:
                return fs
    return min_size


def _uniform_font_size(
    texts: list[str],
    box_width_in: float,
    box_height_in: float | None = None,
    *,
    max_size: int = 24,
    min_size: int = 12,
) -> int:
    """Pick one font size that fits EVERY string in *texts*. Used when a group
    of bullets / step titles / row headers must share the same size for a
    consistent look."""
    if not texts:
        return max_size
    longest = max(texts, key=lambda t: sum(
        (1.0 if ord(c) >= 0x4E00 else 0.55) for c in (t or "")
    ))
    return _dynamic_font_size(
        longest, box_width_in, box_height_in,
        max_size=max_size, min_size=min_size,
    )


def _fit_placeholder_text(
    ph, text: str,
    *,
    max_size: int = 32,
    min_size: int = 18,
    bold: bool | None = None,
    color: str | None = None,
) -> None:
    """Set placeholder text with an adaptively-chosen font size.

    Only size (and optionally bold/color) is overridden — font family, alignment,
    and everything else inherit from the template's placeholder styling.
    """
    if ph is None or not ph.has_text_frame or not text:
        return
    box_w = (ph.width / 914400) if ph.width else 8.0
    box_h = (ph.height / 914400) if ph.height else 1.0
    fs = _dynamic_font_size(text, box_w, box_h, max_size=max_size, min_size=min_size)
    tf = ph.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(fs)
    if bold is not None:
        p.font.bold = bold
    if color:
        p.font.color.rgb = _hex_to_rgb(color)


# ── Per-slide decoration library (logos, watermarks, corner marks) ───

def _extract_decor_library(src_prs) -> dict[str, list[dict]]:
    """Extract per-slide decorative shapes (logos, corner marks, small
    watermarks) from source slides, grouped by role.

    Captures BOTH picture shapes (via image blob) AND arbitrary
    non-picture shapes (via XML deep copy). Many templates author the
    corner logo next to the title as a grouped shape / autoshape /
    freeform with a colour fill rather than an embedded raster image,
    and the previous version silently dropped every non-picture shape.

    Master- and layout-level non-placeholder shapes are inherited
    automatically by add_slide(layout), so this library only needs to
    rescue shapes placed directly on individual source slides — those
    would otherwise be lost when _strip_existing_slides removes the
    source slides. Shapes covering more than 15% of the slide are
    treated as backgrounds and skipped (they belong to the bg library).
    """
    from copy import deepcopy

    slide_area = (src_prs.slide_width or 1) * (src_prs.slide_height or 1)

    def _entries_from(slide) -> list[dict]:
        out: list[dict] = []
        for shape in slide.shapes:
            if shape.is_placeholder:
                continue
            if not shape.width or not shape.height:
                continue
            # Corner logos / marks are almost always small; a hero image
            # or half-slide graphic must not leak into the decor library.
            if shape.width * shape.height > slide_area * 0.1:
                continue
            # Path A: real picture shape. Carry the raw image blob so it
            # can be re-embedded on the destination slide via
            # add_picture — safer than XML-cloning a picture whose rel
            # id does not exist in the new slide's part.
            try:
                blob = shape.image.blob
                ct = shape.image.content_type or ""
                if ct.startswith("image/svg"):
                    # add_picture does not reliably support SVG.
                    continue
                out.append({
                    "kind": "picture",
                    "blob": blob,
                    "content_type": ct,
                    "x": int(shape.left or 0),
                    "y": int(shape.top or 0),
                    "w": int(shape.width),
                    "h": int(shape.height),
                })
                continue
            except (AttributeError, Exception):
                pass
            # Path B: any non-picture decoration — autoshape, group,
            # freeform, text mark, filled rect. Deep-copy the raw XML
            # so the destination slide gets the exact same visual
            # (fills, gradients, text runs, grouped children) without
            # our renderer having to reason about the internal structure.
            #
            # Some templates also place small text labels (corner
            # watermarks, side notes) as freeform text shapes. Those
            # are *content*, not decoration, so we must skip any
            # shape that contains visible text.
            is_text_block = False
            try:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    is_text_block = True
            except Exception:
                pass
            if is_text_block:
                continue
            try:
                out.append({"kind": "xml", "xml": deepcopy(shape._element)})
            except Exception:
                continue
        return out

    lib: dict[str, list[dict]] = {}
    slides = list(src_prs.slides)
    if not slides:
        return lib

    cover = _entries_from(slides[0])
    if cover:
        lib["cover"] = cover
    if len(slides) > 1:
        closing = _entries_from(slides[-1])
        if closing:
            lib["closing"] = closing

    # Content: prefer a middle body slide so we skip cover/closing decorations.
    for s in slides[1:-1] or slides:
        pics = _entries_from(s)
        if pics:
            lib["content"] = pics
            break
    lib.setdefault("content", cover)
    lib.setdefault("section", lib.get("content"))
    return {k: v for k, v in lib.items() if v}


def _send_shape_to_back(shape) -> None:
    """Move *shape* to the bottom of the slide's z-order.

    In OOXML the shape tree is drawn top-down: earlier children render below
    later ones. We remove the shape and reinsert it right after nvGrpSpPr /
    grpSpPr so it sits below any placeholder content.
    """
    el = shape._element
    sp_tree = el.getparent()
    if sp_tree is None:
        return
    sp_tree.remove(el)
    insert_idx = 0
    for i, child in enumerate(sp_tree):
        tag = child.tag.split("}", 1)[-1]
        if tag not in ("nvGrpSpPr", "grpSpPr"):
            insert_idx = i
            break
    else:
        insert_idx = len(sp_tree)
    sp_tree.insert(insert_idx, el)


def _apply_decor(slide, entries: list[dict] | None) -> int:
    """Re-attach captured decorative shapes to *slide*.

    Two paths, matching what _extract_decor_library produced:
      - kind='picture': re-add via slide.shapes.add_picture so the new
        image relationship lives in the destination slide's part.
      - kind='xml':     deep-copy the raw shape XML into the destination
        slide's spTree. This preserves grouped logos, autoshapes,
        freeforms, filled rects, and text marks exactly, without our
        renderer needing to reason about their internal structure.

    Every attached shape is pushed to the back of the z-order so
    placeholder content (title, body) always renders on top.
    """
    if not entries:
        return 0
    import os as _os
    import tempfile
    from copy import deepcopy

    class _ShapeShim:
        # Minimal wrapper so _send_shape_to_back — which only touches
        # ._element and ._element.getparent() — can accept an xml decor
        # entry once we've appended it to the slide's spTree.
        def __init__(self, el):
            self._element = el

    added = 0
    sp_tree = slide.shapes._spTree
    for e in entries:
        kind = e.get("kind") or ("picture" if e.get("blob") else "")
        if kind == "picture" and e.get("blob"):
            ct = e.get("content_type") or ""
            ext = ct.split("/", 1)[1] if "/" in ct else "png"
            fd, path = tempfile.mkstemp(suffix=f".{ext}")
            _os.close(fd)
            try:
                with open(path, "wb") as f:
                    f.write(e["blob"])
                pic = slide.shapes.add_picture(
                    path,
                    Emu(e["x"]), Emu(e["y"]),
                    width=Emu(e["w"]), height=Emu(e["h"]),
                )
                _send_shape_to_back(pic)
                added += 1
            except Exception:
                pass
            finally:
                try:
                    _os.remove(path)
                except OSError:
                    pass
        elif kind == "xml" and e.get("xml") is not None:
            try:
                # deepcopy the source XML so applying the same decor to
                # multiple slides never wires them all to a single node.
                cloned = deepcopy(e["xml"])
                sp_tree.append(cloned)
                _send_shape_to_back(_ShapeShim(cloned))
                added += 1
            except Exception:
                pass
    return added

def _add_slide_footnotes(slide, citations: list[dict], colors) -> None:
    """Render up to 3 citation lines at the slide bottom (small, muted).

    When there are MORE than 3 citations we append a "(+N more, see
    References)" tail so the reader knows where to look. This matches the
    behaviour of the pptxgenjs renderer and keeps the two backends consistent.
    """
    W, H = _slide_dims(slide)
    visible = [c for c in citations if c.get("text")]
    total = len(visible)
    shown = visible[:3]
    lines = [f"{c.get('key', '')} {c.get('text', '')}".strip() for c in shown]
    if total > 3:
        # Small breadcrumb pointing to the full References slide.
        lines.append(f"(+{total - 3} more, see References)")

    base_y = H - 0.55
    line_h = 0.14
    for i, txt in enumerate(lines):
        tb = _new_textbox(
            slide,
            Inches(_MARGIN_L), Inches(base_y + i * line_h),
            Inches(W - _MARGIN_L - _MARGIN_R - 1.3), Inches(line_h),
        )
        p = tb.text_frame.paragraphs[0]
        p.text = (txt[:130] + "…") if len(txt) > 130 else txt
        p.font.size = Pt(9)
        p.font.italic = True
        p.font.color.rgb = _hex_to_rgb(colors.get("muted", "6B7280"))

def _fill_references_slide(slide, s, references: list[dict], colors) -> None:
    """Render the plan-level bibliography as a References slide.

    Prefers the slide's own bullets if the planner already prepared them,
    otherwise builds them from `plan.references` so the deck never ships
    an empty References page.
    """
    items = s.bullets or [f"{r.get('key','')} {r.get('text','')}".strip()
                          for r in references]
    if not items:
        return
    s.bullets = items[:12]
    _fill_body_bullets(slide, s, candidates=(14, 13, 12, 11))
